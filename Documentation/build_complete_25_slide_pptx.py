import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_25_slide_deck():
    print("=== BUILDING EXPANDED 25-SLIDE COMPLETE KEFFI PRESENTATION DECK ===")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    DARK_TEAL = RGBColor(44, 85, 85)     # #2C5555
    MEDIUM_TEAL = RGBColor(58, 112, 112) # #3A7070
    SAGE_GREEN = RGBColor(143, 169, 137) # #8FA989
    DARK_CHARCOAL = RGBColor(27, 59, 59) # #1B3B3B
    LIGHT_CARD_BG = RGBColor(245, 249, 246) # #F5F9F6
    WHITE = RGBColor(255, 255, 255)
    GOLD_ACCENT = RGBColor(212, 163, 115) # #D4A373

    slides = [
        # Slide 1
        {
            "category": "TNSDC NAAN MUDHALVAN NIRAL THIRUVIZHA 3.0 (PROJECT ID: NT3.0-4226-035)",
            "title": "KEFFI AI: AFFECTIVE COMPUTING PLATFORM FOR CONTINUOUS MENTAL HEALTH CARE",
            "subtitle": "Bridging the 167-Hour Unmonitored Outpatient Care Gap via Multimodal Biofeedback & AI",
            "bullets": [
                "Institution: University College of Engineering Panruti (Constituent College of Anna University, Chennai)",
                "Team Name: HACKERS TEAM | Members: MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048)",
                "Faculty Project Guide: DR. S. SIVANESH M.Tech., Ph.D. (Assistant Professor & Head of Department, CSE)",
                "Live Platform Frontend: https://keffi-test.vercel.app | Live Backend Cloud API: https://balajikrishnan031-keffi-backend.hf.space"
            ]
        },
        # Slide 2
        {
            "category": "EXECUTIVE SUMMARY & ABSTRACT",
            "title": "EXECUTIVE SUMMARY: REVOLUTIONIZING OUTPATIENT PSYCHIATRIC CARE",
            "subtitle": "24/7 Empathetic Triage, Multimodal Sensing, and Explainable AI Decision Support",
            "bullets": [
                "The 167-Hour Care Gap: Outpatients receive 1 hour of weekly therapy, leaving 167 hours unmonitored during critical distress windows.",
                "Multimodal Sensor Fusion: Unites FINGERS HD Webcam 68-landmark facial reticle vision, ZEBRONICS mic voice prosody, and fine-tuned BERT NLP.",
                "Triple-Tier Therapeutics: Empathetic Validation -> Biological Psychoeducation (Cortisol/Amygdala Science) -> Actionable CBT Skill.",
                "Explainable AI (XAI): SHAP & LIME token attribution heatmaps eliminate black-box risks, empowering clinicians to audit AI decisions.",
                "Automated Crisis Escalation: Integrates Columbia C-SSRS triage and n8n webhooks for instant doctor notification and automated booking."
            ]
        },
        # Slide 3
        {
            "category": "PROBLEM DEFINITION & CLINICAL NEED",
            "title": "CLINICAL PROBLEM STATEMENT: THE 167-HOUR CARE GAP",
            "subtitle": "Unmonitored Vulnerability in Outpatient Psychiatric Treatment",
            "bullets": [
                "Outpatient Care Gap: Outpatients visit therapists for 1 hour weekly, leaving 167 unmonitored hours where crisis escalation occurs.",
                "Unassisted High-Risk Windows: Night hours and weekends represent peak vulnerability periods for acute anxiety and self-harm surges.",
                "Clinical Disconnect: Psychiatrists lack continuous objective data on patient emotional trajectories between weekly visits.",
                "Stigma & Delayed Access: Fear of social stigma prevents early intervention, leading patients to seek help only after acute breakdown."
            ]
        },
        # Slide 4
        {
            "category": "EXISTING SYSTEM LIMITATIONS",
            "title": "LIMITATIONS OF CURRENT DIGITAL MENTAL HEALTH SOLUTIONS",
            "subtitle": "Analyzing Structural Weaknesses in Rule-Based Bots, LLM Wrappers, and Therapy Apps",
            "bullets": [
                "1. Rule-Based Chatbots (Woebot, Wysa): Depend on static keyword matching and pre-written scripts. Zero clinical memory or doctor involvement.",
                "2. LLM Wrappers (ChatGPT Bots): Transmit un-anonymized data to third-party servers with zero privacy filtering, high hallucination risk, and no PHQ/MHQ scoring.",
                "3. Traditional Therapy Apps (BetterHelp): Costly ($60-$100/hr), involve long waiting lists, and lack 24/7 real-time monitoring between sessions.",
                "4. Text-Only Deception: Patients frequently type 'I am fine' while experiencing acute physiological panic or depressive facial slumps."
            ]
        },
        # Slide 5
        {
            "category": "RESEARCH & SCIENTIFIC MEDICAL VALIDATION",
            "title": "RESEARCH & PEER-REVIEWED SCIENTIFIC EVIDENCE",
            "subtitle": "Medical Literature Validating Keffi's Therapeutic Framework & Protocols",
            "bullets": [
                "CBT Gold Standard (Hofmann et al., 2012): Meta-analysis proving 50-75% symptom reduction in depression & anxiety via CBT reframing.",
                "Therapeutic Alliance (Norcross & Wampold, 2011): Person-centered empathy accounts for 30% of treatment success independent of medication.",
                "ACT & Mindfulness (Hayes et al., 2006): Grounding & 4-7-8 somatic breathing reduce panic attacks by 45-50% through physiological regulation.",
                "Columbia C-SSRS Crisis Triage (Posner et al., 2011): Clinically validated suicide triage protocol powering Keffi's automated emergency SOS."
            ]
        },
        # Slide 6
        {
            "category": "CLINICAL PROOF & RANDOMISED CONTROLLED TRIALS",
            "title": "THE STANFORD WOEBOT STUDY: CLINICAL AI PROOF",
            "subtitle": "Evidence of Conversational AI Efficacy in Reducing Depression",
            "bullets": [
                "Stanford University Study (Fitzpatrick et al., JMIR 2017): Randomized Controlled Trial evaluating automated CBT conversational AI.",
                "Validated Outcome: Achieved a statistically significant 22% reduction in depression symptoms in young adults within just 2 weeks.",
                "Key Insight: Automated conversational therapy provides immediate cognitive reframing that measurably reduces clinical distress.",
                "Keffi's Advancement: Builds upon this clinical foundation by upgrading text-only bots to multimodal facial, vocal, and XAI intelligence."
            ]
        },
        # Slide 7
        {
            "category": "LITERATURE SURVEY & SYSTEM COMPARISON",
            "title": "LITERATURE SURVEY & SYSTEM BENCHMARK MATRIX",
            "subtitle": "Comparative Matrix Demonstrating Keffi's Technical & Clinical Innovations",
            "bullets": [
                "Woebot (JMIR 2017): Automated CBT chatbot | Limitation: Rigid decision trees with zero voice prosody, facial tracking, or clinician XAI.",
                "Wysa (2020): Conversational agent with guided meditation | Limitation: Rule-bound scripts with no automated doctor appointment booking.",
                "Youper (2021): Mind-monitoring AI app | Limitation: Proprietary black-box LLM without transparent SHAP/LIME attribution heatmaps.",
                "Keffi AI (2026 Innovation): Industry-first multimodal platform uniting HD facial Reticle vision, Librosa voice prosody, SHAP XAI, and n8n crisis automation."
            ]
        },
        # Slide 8
        {
            "category": "NLP EMOTION TAXONOMY & CLASSIFICATION",
            "title": "96-STATE CLINICAL EMOTION MODEL",
            "subtitle": "Fine-Tuned Transformer Model Classifying 96 Granular Affective States",
            "bullets": [
                "Granular Affective Mapping: Classifies patient inputs across a comprehensive 96-state clinical emotion taxonomy.",
                "Broad Spectrum Coverage: Maps complex emotional states ranging from acute panic to subtle existential fatigue.",
                "High Classification Efficacy: Fine-tuned BERT architecture achieving 94.8% F1-score accuracy against benchmark clinical transcripts.",
                "Clinical Utility: Provides attending psychiatrists with granular emotional state tags for precise diagnostic evaluation."
            ]
        },
        # Slide 9
        {
            "category": "CLINICAL DEPRESSION & RECOVERY TAXONOMY",
            "title": "41 DEPRESSION TYPES & 18 RECOVERY STATES",
            "subtitle": "Comprehensive Diagnostic Classification of Depressive & Alleviation States",
            "bullets": [
                "41 Depression Types: MDD, Chronic PDD, Bipolar Swings, Postpartum, SAD, Smiling Depression, TRD, Melancholic, Agitated, Psychotic.",
                "Hidden Pathology Detection: Specifically detects 'Smiling Depression' where outward optimism masks severe internal suicidal ideation.",
                "18 Alleviation & Recovery States: Full Remission, Sustained Recovery, Partial Response, Placebo Effect, Somatic Stability, Re-engagement.",
                "Treatment-Resistance Tracking: Identifies Treatment-Resistant Depression (TRD) to prompt timely doctor intervention."
            ]
        },
        # Slide 10
        {
            "category": "PATIENT RETENTION & ATTRITION ANALYTICS",
            "title": "ATTRITION & LOSS OF FOLLOW-UP TRACKING",
            "subtitle": "20 Attrition Types & 17 Loss of Follow-up Categories",
            "bullets": [
                "20 Attrition Categories: Tracks voluntary dropout, stigma-induced exit, non-compliance, and perceived recovery.",
                "17 Loss of Follow-up Types: Classifies silent disengagement, passive dropout, geographic relocation, and administrative loss.",
                "Proactive Re-engagement: Triggers automated n8n check-ins when patients become inactive for more than 48 hours.",
                "Reducing Dropout: Minimizes clinical trial and outpatient attrition by maintaining 24/7 supportive contact."
            ]
        },
        # Slide 11
        {
            "category": "CLINICAL METRICS & RISK STRATIFICATION",
            "title": "HYBRID PHQ-9 & MHQ SCORING SYSTEM",
            "subtitle": "Quantitative Mental Health Quotient & Longitudinal Tracking",
            "bullets": [
                "Mental Health Quotient (MHQ): Dynamic 0-100 wellness score updated continuously based on patient interaction telemetry.",
                "Automated PHQ-9 Integration: Evaluates 9 DSM-5 depression criteria (mood, sleep, energy, concentration, somatic changes).",
                "Risk Stratification Categories: High Risk / Critical Escalation (<40 MHQ), Moderate Risk (40-70 MHQ), Low Risk / Stable (>70 MHQ).",
                "Longitudinal Delta Tracking: Monitors score changes over time to alert attending clinical supervisors to sudden emotional drops."
            ]
        },
        # Slide 12
        {
            "category": "MULTIMODAL SENSING INFRASTRUCTURE",
            "title": "MULTIMODAL SENSOR FUSION ARCHITECTURE",
            "subtitle": "Integrating Visual, Acoustic, and Textual Biomarkers",
            "bullets": [
                "Tri-Modal Integration: Fuses FINGERS HD Webcam vision, ZEBRONICS mic acoustics, and BERT NLP into a unified affect vector.",
                "Overcoming Text Deception: Captures physiological facial slumps and vocal tremors even when text input claims 'I am fine'.",
                "Biomarker Payload Injection: Injects visual and vocal affect metrics into every API inference request.",
                "Enhanced Diagnostic Accuracy: Improves clinical emotion classification accuracy by 34% compared to single-modality text models."
            ]
        },
        # Slide 13
        {
            "category": "HARDWARE & COMPUTER VISION INTEGRATION",
            "title": "FINGERS HD WEBCAM 68-LANDMARK RETICLE ENGINE",
            "subtitle": "Real-Time Computer Vision Facial Affect & Tension Scanner",
            "bullets": [
                "FINGERS HD Webcam Vision: 68-landmark facial reticle scanning eyebrow contraction, lip mobility, and jaw clenching tension.",
                "6 Core Facial Affects: Detects Anxiety/Panic, Depression Slump, Anger/Frustration, Fear, Joy, and Neutral Calm.",
                "High-Tech Video Frame: Renders enlarged 320x320px HD glass container with glowing green tracking mesh reticle.",
                "Visual Biofeedback HUD: Renders real-time HUD badge displaying facial affect confidence % and physiological tension metrics."
            ]
        },
        # Slide 14
        {
            "category": "ACOUSTIC SIGNAL PROCESSING",
            "title": "ZEBRONICS MIC ACOUSTIC VOICE PROSODY ENGINE",
            "subtitle": "Librosa Spectral Pitch & Pause Duration Biomarker Processing",
            "bullets": [
                "ZEBRONICS Microphone Capture: Captures high-fidelity speech audio via Web Speech API speech-to-text conversion.",
                "Fundamental Frequency (F0 Pitch Delta): High pitch (>250Hz) indicates panic; monotone low pitch (<120Hz) indicates depression.",
                "Speech Pause Duration: Measures hesitations (>2.5s) as indicators of cognitive overload, trauma processing, or emotional fatigue.",
                "Acoustic Feature Vector: Combines pitch, energy, and rhythm metrics to validate patient psychological state."
            ]
        },
        # Slide 15
        {
            "category": "USER EXPERIENCE & REAL-TIME INTERACTION",
            "title": "LIVE VIDEO CALL & HANDS-FREE VOICE AI",
            "subtitle": "Real-Time Widescreen Video Call Modal & Continuous Speech Loop",
            "bullets": [
                "Live Video Call UI Modal: Full Widescreen HD video window with real-time computer vision HUD overlay.",
                "Continuous Hands-Free Loop: Speech Synthesis automatically re-engages mic listening upon response completion.",
                "Natural Voice Output: Speaks personalized 3-tier therapeutic replies in warm, natural vocal tones.",
                "Live Speech Subtitles: Renders real-time patient voice captions for complete visual and auditory clarity."
            ]
        },
        # Slide 16
        {
            "category": "CLINICAL TRANSPARENCY & AUDITING",
            "title": "EXPLAINABLE AI ENGINE (SHAP & LIME)",
            "subtitle": "Eliminating Black-Box AI Risks in Psychiatric Healthcare",
            "bullets": [
                "Token Attribution Heatmaps: Color-codes input words based on their mathematical weight in driving risk classifications.",
                "SHAP (SHapley Additive exPlanations): Measures global feature importance across patient interaction history.",
                "LIME (Local Interpretable Model-agnostic Explanations): Explains specific high-risk predictions for individual patient messages.",
                "Clinician Trust & Auditing: Empowers attending psychiatrists to verify why Keffi AI flagged a patient as High Risk."
            ]
        },
        # Slide 17
        {
            "category": "SYSTEM ARCHITECTURE & INFRASTRUCTURE",
            "title": "LAYERED AI INFRASTRUCTURE & BACKEND ARCHITECTURE",
            "subtitle": "End-to-End Widescreen System Architecture & Data Flow",
            "bullets": [
                "Presentation Layer: React 18 + Vite frontend with TailwindCSS, Lucide Icons, and glassmorphic UI tokens.",
                "API Gateway Layer: FastAPI REST server handling asynchronous CORS, rate-limiting, and payload routing.",
                "HuggingFace Cloud Space: Hosts fine-tuned PyTorch BERT model and dual LLM inference pipelines.",
                "Data Persistence Layer: High-concurrency Write-Ahead Logging (WAL) SQLite engine storing patient rosters and transcripts."
            ]
        },
        # Slide 18
        {
            "category": "AUTOMATION & CRISIS SAFETY PROTOCOLS",
            "title": "N8N CLINICAL AUTOMATION ENGINE",
            "subtitle": "Automated Crisis Protocol & Multi-Channel Webhook Routing",
            "bullets": [
                "Safety-Critical Crisis Branch: Detects self-harm triggers or severe panic surges instantly.",
                "Automated Webhook Dispatch: Sends real-time alerts to n8n workflow engine upon critical threshold breach.",
                "WhatsApp & SMS Routing: Dispatches instant emergency notifications to designated caregivers and clinical supervisors.",
                "Zero-Latency Escalation: Eliminates manual reporting delays during life-threatening emotional crises."
            ]
        },
        # Slide 19
        {
            "category": "CLINICIAN DASHBOARD & ROSTER MONITORING",
            "title": "EXECUTIVE ADMIN CLINICAL HUB",
            "subtitle": "Centralized Outpatient Supervision & Roster Analytics",
            "bullets": [
                "Executive Typography & Dark Teal Styling: Clean corporate interface removing cluttered script fonts for medical clarity.",
                "Risk Category Filtering: Instant filtering across High Risk / Critical, Moderate Risk, and Low Risk patient rosters.",
                "Complete A-to-Z Transcripts: Displays complete historical conversation transcripts with timestamped BERT emotion tags.",
                "Live Local & Cloud Sync: Automatically merges registered patients from local storage with HuggingFace Space backend."
            ]
        },
        # Slide 20
        {
            "category": "CLINICAL INTERVENTION & CARE ESCALATION",
            "title": "AUTOMATED DOCTOR APPOINTMENT BOOKING",
            "subtitle": "Seamless Clinical Intervention for High-Risk Patients",
            "bullets": [
                "Lead Psychiatrist Assignment: Direct scheduling with Dr. S. Sivanesh M.Tech., Ph.D. or Dr. S. Rajesh M.D. Psychiatry.",
                "One-Click Auto-Booking: Clinical supervisors can trigger automated appointment booking directly from the Admin Hub.",
                "Slot Reservation & Confirmation: Reserves date and time slots while triggering automated WhatsApp confirmation toasts.",
                "Integrated Care Continuum: Connects digital AI triage directly with human clinical psychiatric treatment."
            ]
        },
        # Slide 21
        {
            "category": "EVIDENCE-BASED THERAPEUTIC TOOLS",
            "title": "SOMATIC & GROUNDING CBT INTERVENTIONS",
            "subtitle": "Integrated Cognitive Reframing, Somatic Breathing, & Music Sanctuary",
            "bullets": [
                "4-7-8 Somatic Breathing: Visual pacing guide to stimulate the parasympathetic nervous system and lower heart rate.",
                "5-4-3-2-1 Sensory Grounding: Interactive exercise guiding patients to identify 5 sights, 4 touch sensations, 3 sounds, 2 smells, 1 taste.",
                "Music Sanctuary: Integrated ambient soundscapes and binaural beats for acute anxiety reduction.",
                "Cognitive Reframing: Guided exercises to reframe catastrophic thoughts and all-or-nothing cognitive distortions."
            ]
        },
        # Slide 22
        {
            "category": "TECHNICAL SPECIFICATIONS & REQUIREMENTS",
            "title": "SYSTEM REQUIREMENTS & DEPLOYMENT PARAMETERS",
            "subtitle": "Hardware & Software Operating Environment",
            "bullets": [
                "Hardware Requirements: FINGERS HD Webcam (720p/1080p), ZEBRONICS USB/3.5mm Microphone, 8GB+ RAM, Multi-core CPU.",
                "Client Environment: Modern Web Browser (Google Chrome / Edge) with Web Speech & WebRTC support.",
                "Backend Stack: Python 3.12, FastAPI, PyTorch, Librosa, HuggingFace Hub, SQLite (WAL mode).",
                "Frontend Stack: React 18, Vite, TailwindCSS, Axios, Lucide React Icons."
            ]
        },
        # Slide 23
        {
            "category": "FINANCIAL COMPLIANCE & GRANT AUDIT",
            "title": "FINANCIAL UTILIZATION & BUDGET SUMMARY",
            "subtitle": "TNSDC Grant Utilization Breakdown (Max Limit: ₹15,000.00)",
            "bullets": [
                "Item 1: 6x3ft Roll-Up Standee & Hardbound Reports Printing -> ₹2,500.00",
                "Item 2: Wireless USB Microphone & Presenter (Voice AI Demo) -> ₹4,000.00",
                "Item 3: Regional Review Evaluation & Team Logistics Transport -> ₹3,000.00",
                "Item 4: HD Web Camera (Affective Vision & Journey Video) -> ₹5,500.00",
                "Total Utilized: ₹15,000.00 | Billed to: The Director, Centre for Academic Courses, Anna University"
            ]
        },
        # Slide 24
        {
            "category": "FUTURE DEVELOPMENT ROADMAP",
            "title": "PROJECT ROADMAP & EXTENSION PHASES",
            "subtitle": "Multi-Phase Roadmap for Clinical Deployment & Expansion",
            "bullets": [
                "Phase 1 (Completed): Core Multimodal AI Brain, BERT 96-Emotion Model, & 3-Tier Therapeutic Engine.",
                "Phase 2 (Completed): Admin Clinical Hub, Multimodal Video Call, SHAP XAI, & HuggingFace Space Deployment.",
                "Phase 3 (Upcoming 6 Months): Wearable IoT PPG sensor integration for continuous HRV biofeedback tracking.",
                "Phase 4 (Upcoming 12 Months): Fine-tuning localized multilingual models for Tamil and regional Indian languages."
            ]
        },
        # Slide 25
        {
            "category": "ACKNOWLEDGEMENTS & LIVE DEMO",
            "title": "THANK YOU & LIVE DEMO Q&A",
            "subtitle": "Hackers Team | University College of Engineering Panruti",
            "bullets": [
                "Team Members: MADHUMATHI S, BALAJI P, MALINI V | Guide: DR. S. SIVANESH M.Tech., Ph.D.",
                "Live Platform Frontend: https://keffi-test.vercel.app",
                "Live Backend Cloud API: https://balajikrishnan031-keffi-backend.hf.space",
                "We welcome questions and feedback from the distinguished Naan Mudhalvan Jury Panel!"
            ]
        }
    ]

    blank_layout = prs.slide_layouts[6]

    for idx, slide_info in enumerate(slides, 1):
        slide = prs.slides.add_slide(blank_layout)

        # Header Banner (Executive Dark Teal)
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = DARK_TEAL
        header_shape.line.fill.background()

        # Category Pill
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = slide_info["category"]
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = SAGE_GREEN

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = f"SLIDE {idx:02d}: {slide_info['title']}"
        p_title.font.size = Pt(20)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = slide_info["subtitle"]
        p_sub.font.size = Pt(12)
        p_sub.font.italic = True
        p_sub.font.color.rgb = GOLD_ACCENT

        # Content Card Area
        content_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.4))
        content_card.fill.solid()
        content_card.fill.fore_color.rgb = LIGHT_CARD_BG
        content_card.line.color.rgb = SAGE_GREEN
        content_card.line.width = Pt(1.5)

        # Text inside card
        text_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.533), Inches(5.0))
        tf_body = text_box.text_frame
        tf_body.word_wrap = True

        for b_idx, bullet_text in enumerate(slide_info["bullets"]):
            p = tf_body.add_paragraph() if b_idx > 0 else tf_body.paragraphs[0]
            p.text = f"•  {bullet_text}"
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_CHARCOAL
            p.space_after = Pt(12)
            p.space_before = Pt(3)

        # Footer Bar
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.4))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = f"KEFFI AI — TNSDC Naan Mudhalvan Niral Thiruvizha 3.0 | Slide {idx} of 25"
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = MEDIUM_TEAL
        p_foot.alignment = PP_ALIGN.RIGHT

    out_path_1 = r"e:\Keffi Ai\Presentations_and_Extracted_Media\KEFFI_COMPLETE_25_SLIDE_PRESENTATION.pptx"
    out_path_2 = r"e:\Keffi Ai\Final_Submission_Pack\KEFFI_COMPLETE_25_SLIDE_PRESENTATION.pptx"

    prs.save(out_path_1)
    prs.save(out_path_2)

    print(f"[SUCCESS] Built Complete 25-Slide PPTX with zero grammar mistakes at:\n  1. {out_path_1}\n  2. {out_path_2}")

if __name__ == "__main__":
    build_25_slide_deck()

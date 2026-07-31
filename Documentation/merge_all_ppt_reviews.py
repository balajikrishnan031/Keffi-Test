"""
Master PPT Merger for Keffi AI - Niral Thiruvizha 3.0 Final Review
Merges slides from FIRST REVIEW.pptx & SECOND REVIEW-1.pptx and adds all Keffi AI upgrades.
Saves output to Presentations_and_Extracted_Media/FINAL_REVIEW_KEFFI_AI_MASTER_COMBINED.pptx
"""

import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

PPT1_PATH = r"e:\Keffi Ai\Presentations_and_Extracted_Media\FIRST REVIEW.pptx"
PPT2_PATH = r"e:\Keffi Ai\Presentations_and_Extracted_Media\SECOND REVIEW-1.pptx"
OUTPUT_PPTX = r"e:\Keffi Ai\Presentations_and_Extracted_Media\FINAL_REVIEW_KEFFI_AI_MASTER_COMBINED.pptx"

def create_master_combined_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Widescreen
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    COLOR_TEAL_DARK = RGBColor(13, 80, 80)     # #0D5050
    COLOR_TEAL_MID = RGBColor(44, 85, 85)     # #2C5555
    COLOR_SLATE_DARK = RGBColor(30, 41, 59)   # #1E293B
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_LIGHT_BG = RGBColor(241, 245, 249)  # #F1F5F9
    COLOR_ACCENT = RGBColor(42, 168, 112)     # #2AA870

    def add_header(slide, title_text, category="NIRAL THIRUVIZHA 3.0 - MASTER FINAL REVIEW DECK"):
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_TEAL_DARK
        header.line.fill.background()

        tf = header.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.margin_top = Inches(0.15)

        p_sub = tf.paragraphs[0]
        p_sub.text = category.upper()
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_ACCENT

        p_main = tf.add_paragraph()
        p_main.text = title_text
        p_main.font.size = Pt(20)
        p_main.font.bold = True
        p_main.font.color.rgb = COLOR_WHITE

    # Extract text content from PPT 1 and PPT 2
    slides_data = []
    for ppt_label, ppt_path in [("REVIEW 1", PPT1_PATH), ("REVIEW 2", PPT2_PATH)]:
        if os.path.exists(ppt_path):
            source_prs = Presentation(ppt_path)
            for idx, slide in enumerate(source_prs.slides, 1):
                slide_info = {"source": ppt_label, "idx": idx, "texts": []}
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        if t:
                            slide_info["texts"].append(t)
                slides_data.append(slide_info)

    # 1. MASTER TITLE SLIDE
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_TEAL_DARK
    bg1.line.fill.background()

    tf1 = bg1.text_frame
    tf1.margin_left = Inches(1.0)
    tf1.margin_top = Inches(1.0)

    p = tf1.paragraphs[0]
    p.text = "NIRAL THIRUVIZHA 3.0 (2026) — COMBINED MASTER FINAL REVIEW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    p = tf1.add_paragraph()
    p.text = "KEFFI AI: Master Clinical AI Psychiatrist & Affective Computing Engine"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tf1.add_paragraph()
    p.text = "Comprehensive Combined Presentation (First Review + Second Review + Clinical Upgrades)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(226, 232, 240)

    box1 = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.5), Inches(11.333), Inches(3.4))
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_TEAL_MID
    box1.line.color.rgb = COLOR_ACCENT

    tf_b1 = box1.text_frame
    tf_b1.margin_left = Inches(0.4)
    tf_b1.margin_top = Inches(0.25)

    info_list = [
        ("Team ID:", "NMNTSTD42260064   |   Team Name: Hackers Team"),
        ("College Name:", "University College of Engineering, Panruti (UCE, Panruti)"),
        ("Problem Statement:", "How might we utilize AI chatbots and machine learning to address incomplete alleviation of depression symptoms, attrition, and loss of follow-up in mental health treatment"),
        ("Theme:", "Artificial Intelligence (Healthcare & Clinical Digital Therapeutics)"),
        ("Team Leader:", "Mathu (Madhumathi S - Reg: 422623104003)"),
        ("Team Members:", "Balaji P (Reg: 4226231035), Malini V (Reg: 4226231048)"),
        ("Faculty Guide:", "Dr. R. Sivanesh, M.E., Ph.D. (Associate Professor, Dept. of CSE)"),
        ("Review Schedule:", "05/08/2026 Afternoon Session (AN) | Panel 2 | Chennai Region Venue")
    ]
    for lbl, val in info_list:
        p = tf_b1.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{lbl} "
        r1.font.bold = True
        r1.font.size = Pt(11.5)
        r1.font.color.rgb = COLOR_ACCENT
        
        r2 = p.add_run()
        r2.text = val
        r2.font.size = Pt(11.5)
        r2.font.color.rgb = COLOR_WHITE

    # COMBINE SLIDES FROM REVIEW 1 & REVIEW 2
    for item in slides_data:
        s = prs.slides.add_slide(blank_layout)
        header_title = f"{item['source']} — SLIDE {item['idx']}"
        if item['texts']:
            first_line = item['texts'][0].replace('\n', ' ')
            if len(first_line) > 60:
                first_line = first_line[:60] + "..."
            header_title = f"{item['source']} (Slide {item['idx']}): {first_line}"

        add_header(s, header_title)

        c_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = COLOR_LIGHT_BG
        c_box.line.color.rgb = COLOR_TEAL_DARK

        tf = c_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.4)
        tf.margin_top = Inches(0.3)

        for text_block in item['texts']:
            p = tf.add_paragraph()
            p.text = text_block
            p.font.size = Pt(12)
            p.font.color.rgb = COLOR_SLATE_DARK

    # APPEND ADVANCED KEFFI CLINICAL UPGRADES SLIDES
    upgrade_slides = [
        ("UPGRADE 1: 4 NEXT-GEN AFFECTIVE COMPUTING LAYERS", [
            "Layer 1: Voice Prosody & Acoustic Sentiment Analyzer (Pitch, Speech Rate WPM, Energy dB, Pause Gaps).",
            "Layer 2: Cognitive Distortion Mapping Engine (10 CBT psychological distortions: Catastrophizing, All-or-Nothing, Mind Reading).",
            "Layer 3: IoT Biometric Telemetry Engine (ESP32 Heart Rate BPM, HRV ms, GSR uS sensor streams for HR > 100 BPM panic override).",
            "Layer 4: Temporal Emotion Knowledge Graph (Long-term graph memory: User -> Triggers -> States -> Coping)."
        ]),
        ("UPGRADE 2: MASTER CLINICAL KNOWLEDGE BASE (DSM-5-TR)", [
            "Major Depressive Disorder (296.2x / 6A70): First-line Cognitive Therapy (Beck) + Behavioral Activation.",
            "Persistent Depressive Disorder / Dysthymia (300.4): Long-term ACT Defusion (Hayes) + Core Schema Reframing.",
            "Generalized Anxiety Disorder (300.02 / 6B00): Somatic Nervous System Pacing (van der Kolk) + Polyvagal 4-7-8 Breathing.",
            "Panic Disorder & Agoraphobia (300.01): TIPP Crisis Distress Tolerance (Linehan DBT) + Interoceptive Exposure.",
            "PTSD & Trauma (309.81 / 6B40): Somatic Neurobiology & Polyvagal Vagus Nerve Pacing."
        ]),
        ("UPGRADE 3: EXPLAINABLE AI (SHAP & LIME XAI SUITE)", [
            "SHAP (SHapley Additive exPlanations): Game-theoretic feature attribution for every input token driving BERT emotion scores.",
            "LIME (Local Interpretable Model-agnostic Explanations): Local surrogate model explaining why CBT reframing was chosen.",
            "100% Medical Audit Compliance: Complete transparency eliminating AI black box in clinical decision making."
        ]),
        ("UPGRADE 4: 70B MULTI-LLM RESILIENT CASCADE", [
            "Primary Engine: Groq Llama-3.3 70B (llama-3.3-70b-versatile) — High-speed clinical response inference.",
            "Fallback 1: OpenAI ChatGPT (gpt-4o-mini) — Failover for complex conversational reframing.",
            "Fallback 2: Google Gemini (gemini-1.5-flash) — Context analysis failover.",
            "Fallback 3: Local 10-Methodology CBT Engine — Zero-downtime offline clinical fallback."
        ])
    ]

    for title, bullets in upgrade_slides:
        s = prs.slides.add_slide(blank_layout)
        add_header(s, title)

        c_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
        c_box.fill.solid()
        c_box.fill.fore_color.rgb = COLOR_TEAL_MID
        c_box.line.color.rgb = COLOR_ACCENT

        tf = c_box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.4)
        tf.margin_top = Inches(0.4)

        for b in bullets:
            p = tf.add_paragraph()
            p.text = f"•  {b}"
            p.font.size = Pt(13.5)
            p.font.color.rgb = COLOR_WHITE

    prs.save(OUTPUT_PPTX)
    print(f"  [SUCCESS] Created Master Combined PPT: {OUTPUT_PPTX} ({len(prs.slides)} slides)")

if __name__ == "__main__":
    create_master_combined_ppt()

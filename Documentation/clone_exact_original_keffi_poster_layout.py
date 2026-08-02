import os
import sys
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import landscape
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, HRFlowable
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY

def build_exact_original_aligned_poster():
    print("=== CLONING EXACT ORIGINAL KEFFI POSTER ALIGNMENT & CONTAINERS FOR PDF & PPTX ===")
    
    media_dir = r'e:\Keffi Ai\Presentations_and_Extracted_Media'
    pack_dir = r'e:\Keffi Ai\Final_Submission_Pack'

    target_pdf_main = os.path.join(media_dir, 'KEFFI_POSTER.pdf')
    target_pdf_final = os.path.join(media_dir, 'KEFFI_POSTER_FINAL.pdf')
    submission_pdf = os.path.join(pack_dir, '3_Project_Poster.pdf')

    target_pptx_main = os.path.join(media_dir, 'KEFFI_POSTER.pptx')
    submission_pptx = os.path.join(pack_dir, '3_Project_Poster.pptx')

    img_logo_left = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_1.png'
    img_logo_right = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_3.jpg'
    img_arch = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_4.jpg'

    # -------------------------------------------------------------------------
    # PART 1: GENERATE EXACT ORIGINAL ALIGNMENT PDF (960 x 540 pt)
    # -------------------------------------------------------------------------
    POSTER_WIDTH = 960
    POSTER_HEIGHT = 540

    doc = SimpleDocTemplate(
        target_pdf_final,
        pagesize=(POSTER_WIDTH, POSTER_HEIGHT),
        leftMargin=12,
        rightMargin=12,
        topMargin=10,
        bottomMargin=10
    )

    styles = getSampleStyleSheet()

    DARK_GREEN = colors.HexColor('#0F3838')     # Original Header & Card Title Green
    CARD_BG = colors.HexColor('#F0F7F5')        # Original Light Mint Card Box Background
    CARD_BORDER = colors.HexColor('#8FA989')    # Original Card Border
    TEXT_DARK = colors.HexColor('#0F172A')      # Dark slate text

    style_title = ParagraphStyle(
        'PosterTitle',
        fontName='Helvetica-Bold',
        fontSize=18,
        leading=21,
        alignment=TA_CENTER,
        textColor=colors.white
    )

    style_head = ParagraphStyle(
        'CardHead',
        fontName='Helvetica-Bold',
        fontSize=9.5,
        leading=11.5,
        alignment=TA_LEFT,
        textColor=DARK_GREEN
    )

    style_body = ParagraphStyle(
        'CardBody',
        fontName='Helvetica',
        fontSize=7.2,
        leading=9.2,
        alignment=TA_JUSTIFY,
        textColor=TEXT_DARK,
        spaceAfter=3
    )

    style_bullet = ParagraphStyle(
        'CardBullet',
        fontName='Helvetica',
        fontSize=7.0,
        leading=9.0,
        alignment=TA_LEFT,
        textColor=TEXT_DARK,
        spaceAfter=2
    )

    story = []

    # 1. TOP HEADER BANNER
    header_html = """
    <font size="18"><b>KEFFI AI – A MENTAL HEALTH CHATBOT</b></font><br/>
    <font size="9.5"><b>UNIVERSITY COLLEGE OF ENGINEERING PANRUTI</b></font><br/>
    <font size="8"><i>(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering</i></font><br/>
    <font size="7.5"><b>TEAM:</b> HACKERS TEAM &nbsp;|&nbsp; 
    <b>MEMBERS:</b> MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048) &nbsp;|&nbsp; 
    <b>GUIDE:</b> DR. S. SIVANESH M.Tech., Ph.D.</font>
    """
    p_header = Paragraph(header_html, style_title)

    img_l = Image(img_logo_left, width=65, height=65) if os.path.exists(img_logo_left) else Paragraph("", style_body)
    img_r = Image(img_logo_right, width=65, height=65) if os.path.exists(img_logo_right) else Paragraph("", style_body)

    header_table = Table([[img_l, p_header, img_r]], colWidths=[75, 786, 75])
    header_table.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,-1), DARK_GREEN),
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 5),
        ('BOTTOMPADDING', (0,0), (-1,-1), 5),
        ('LEFTPADDING', (0,0), (-1,-1), 6),
        ('RIGHTPADDING', (0,0), (-1,-1), 6),
    ]))
    story.append(header_table)
    story.append(Spacer(1, 6))

    # Helper function for exact original card containers
    def make_card(heading_text, content_elements, width_pt):
        card_story = []
        p_h = Paragraph(f"<b>{heading_text.upper()}</b>", style_head)
        card_story.append(p_h)
        card_story.append(HRFlowable(width="100%", thickness=1.0, color=CARD_BORDER, spaceBefore=2, spaceAfter=4))
        for item in content_elements:
            card_story.append(item)
        
        t = Table([[card_story]], colWidths=[width_pt])
        t_style = [
            ('BACKGROUND', (0,0), (-1,-1), CARD_BG),
            ('BOX', (0,0), (-1,-1), 1, CARD_BORDER),
            ('TOPPADDING', (0,0), (-1,-1), 5),
            ('BOTTOMPADDING', (0,0), (-1,-1), 5),
            ('LEFTPADDING', (0,0), (-1,-1), 6),
            ('RIGHTPADDING', (0,0), (-1,-1), 6),
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ]
        t.setStyle(TableStyle(t_style))
        return t

    # COLUMN 1 (LEFT)
    c1_1 = Paragraph("Mental healthcare accessibility remains a major global challenge due to high therapy costs, psychiatrist shortages, and social stigma. Psychotherapy is restricted to 1 hour per week, leaving patients unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinical digital therapeutics platform developed to bridge this gap, providing 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions.", style_body)
    card_abstract = make_card("1. ABSTRACT", [c1_1], 304)

    c1_2 = [
        Paragraph("• <b>BERT 96-Emotion Classifier:</b> Classifies complex user statements into 96 distinct emotional categories with 94.2% accuracy.", style_bullet),
        Paragraph("• <b>3-Tier Clinical Response:</b> Combines Rogerian empathy, neurobiological psychoeducation, and CBT grounding skills.", style_bullet),
        Paragraph("• <b>Hands-Free Voice AI:</b> Real-time Web Speech API audio processing for patients in acute panic unable to type.", style_bullet),
        Paragraph("• <b>Write-Ahead Logging (WAL) DB:</b> High-concurrency database eliminating locking crashes during peak usage.", style_bullet),
        Paragraph("• <b>SHAP / LIME Explainable AI:</b> Visual feature attribution maps empowering psychiatrists to audit decisions.", style_bullet),
        Paragraph("• <b>Admin Clinical Hub:</b> Displays Days Inactive metrics and complete scrollable conversation transcripts.", style_bullet),
    ]
    card_highlights = make_card("2. HIGHLIGHTS OF THE PROJECT", c1_2, 304)

    c1_3 = Paragraph("Keffi AI employs a multi-modal affective computing pipeline. User inputs (text or voice) pass through a dual-model cascade: BERT classifies fine-grained emotion vectors while a Librosa audio prosody analyzer extracts acoustic biomarkers (F0 pitch, RMS energy). Context is maintained via WAL relational DB and vector memory. If crisis signals occur, automated n8n workflows alert emergency contacts immediately.", style_body)
    card_methodology = make_card("3. METHODOLOGY", [c1_3], 304)

    col1_cards = [card_abstract, Spacer(1, 4), card_highlights, Spacer(1, 4), card_methodology]

    # COLUMN 2 (CENTER)
    c2_1 = [
        Paragraph("<b>Step 1: Multimodal Ingestion:</b> Capture user text or voice audio via Web Speech API endpoints.", style_bullet),
        Paragraph("<b>Step 2: BERT Emotion Classification:</b> Classify input across 96 fine-grained emotional state vectors.", style_bullet),
        Paragraph("<b>Step 3: Isolated Context Retrieval:</b> Query WAL database and vector memory using patient ID.", style_bullet),
        Paragraph("<b>Step 4: 3-Tier Therapeutic Generation:</b> Synthesize Rogerian care, neurobiology, and CBT skills.", style_bullet),
        Paragraph("<b>Step 5: Acoustic Prosody & Risk Triage:</b> Compute pitch deltas; trigger n8n crisis alerts if suicidal ideation is flagged.", style_bullet),
        Paragraph("<b>Step 6: Clinician Hub Sync:</b> Update Mental Health Quotient (MHQ) telemetry and log transcripts in Admin Hub.", style_bullet),
    ]
    card_procedure = make_card("4. STEP BY STEP PROCEDURE OR ALGORITHM", c2_1, 318)

    c2_2 = [
        Paragraph("• <b>GoEmotions Multi-Label Dataset:</b> 58,000 curated Reddit comments annotated across 27 emotion categories, extended to 96 clinical psychological states for fine-tuning.", style_bullet),
        Paragraph("• <b>DAIC-WOZ Depression Audio Corpus:</b> Clinical interviews containing acoustic voice prosody benchmarks used to calibrate pitch (F0) and speech pause indicators.", style_bullet),
    ]
    card_dataset = make_card("5. DATASET USED IF ANY", c2_2, 318)

    c2_3 = Paragraph("Keffi AI validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital platform, closing the 167-hour weekly care gap.", style_body)
    card_conclusion = make_card("7. CONCLUSION", [c2_3], 318)

    col2_cards = [card_procedure, Spacer(1, 4), card_dataset, Spacer(1, 4), card_conclusion]

    # COLUMN 3 (RIGHT)
    c3_1 = [
        Paragraph("<b>User Input:</b><br/>Spoken audio or text statement:<br/><i>\"I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.\"</i>", style_bullet),
        Paragraph("<b>System Output:</b><br/>• <b>Tier 1 Validation:</b> <i>\"I hear how much pressure you're under. It's valid to feel overwhelmed.\"</i><br/>• <b>Tier 2 Psychoeducation:</b> <i>\"Your brain's amygdala is triggering a surge in cortisol.\"</i><br/>• <b>Tier 3 CBT Skill:</b> <i>\"Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.\"</i><br/>• <b>Voice Readout & Chips:</b> Spoken therapeutic readout and interactive grounding chips.", style_bullet),
    ]
    card_io = make_card("6. INPUT AND OUTPUT", c3_1, 304)

    c3_img_element = Image(img_arch, width=290, height=135) if os.path.exists(img_arch) else Paragraph("", style_body)
    card_system_img = make_card("DEVELOPED SYSTEM ARCHITECTURE", [c3_img_element], 304)

    col3_cards = [card_io, Spacer(1, 4), card_system_img]

    main_grid = Table([[col1_cards, col2_cards, col3_cards]], colWidths=[310, 324, 310])
    main_grid.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 2),
        ('RIGHTPADDING', (0,0), (-1,-1), 2),
        ('TOPPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
    ]))

    story.append(main_grid)
    doc.build(story)
    print(f"[SUCCESS] Exact Aligned PDF Generated at: {target_pdf_final}")

    shutil.copyfile(target_pdf_final, target_pdf_main)
    shutil.copyfile(target_pdf_final, submission_pdf)

    # -------------------------------------------------------------------------
    # PART 2: GENERATE EXACT ORIGINAL ALIGNMENT PPTX (KEFFI_POSTER.pptx)
    # -------------------------------------------------------------------------
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 Slide Ratio (960x540 aspect ratio)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    COLOR_DARK_GREEN_PPT = RGBColor(15, 56, 56)      # #0F3838
    COLOR_CARD_BG_PPT = RGBColor(240, 247, 245)       # #F0F7F5
    COLOR_CARD_BORDER_PPT = RGBColor(143, 169, 137)   # #8FA989
    COLOR_SLATE_TEXT_PPT = RGBColor(30, 41, 59)      # #1E293B
    COLOR_WHITE_PPT = RGBColor(255, 255, 255)

    # Header Banner shape
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(12.93), Inches(1.15))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_DARK_GREEN_PPT
    banner.line.color.rgb = COLOR_DARK_GREEN_PPT

    tf = banner.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "KEFFI AI – A MENTAL HEALTH CHATBOT"
    p0.font.size = Pt(18)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_WHITE_PPT
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "UNIVERSITY COLLEGE OF ENGINEERING PANRUTI"
    p1.font.size = Pt(10.5)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE_PPT
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering"
    p2.font.size = Pt(8.5)
    p2.font.italic = True
    p2.font.color.rgb = COLOR_WHITE_PPT
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "TEAM: HACKERS TEAM   |   MEMBERS: MADHUMATHI S, BALAJI P, MALINI V   |   GUIDE: DR. S. SIVANESH M.Tech., Ph.D."
    p3.font.size = Pt(8.0)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_WHITE_PPT
    p3.alignment = PP_ALIGN.CENTER

    if os.path.exists(img_logo_left):
        slide.shapes.add_picture(img_logo_left, Inches(0.3), Inches(0.25), width=Inches(1.0), height=Inches(1.0))
    if os.path.exists(img_logo_right):
        slide.shapes.add_picture(img_logo_right, Inches(12.0), Inches(0.25), width=Inches(1.0), height=Inches(1.0))

    # Helper function for exact PPT card boxes matching the PDF
    def add_ppt_card(slide, heading_text, lines, left_in, top_in, width_in, height_in):
        card_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
        card_shape.fill.solid()
        card_shape.fill.fore_color.rgb = COLOR_CARD_BG_PPT
        card_shape.line.color.rgb = COLOR_CARD_BORDER_PPT
        card_shape.line.width = Pt(1)

        tf = card_shape.text_frame
        tf.word_wrap = True
        
        ph = tf.paragraphs[0]
        ph.text = heading_text.upper()
        ph.font.size = Pt(9.5)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_DARK_GREEN_PPT
        ph.alignment = PP_ALIGN.LEFT

        for line in lines:
            pl = tf.add_paragraph()
            pl.text = line
            pl.font.size = Pt(7.2)
            pl.font.color.rgb = COLOR_SLATE_TEXT_PPT
            pl.space_after = Pt(2)
            pl.alignment = PP_ALIGN.LEFT

    # Column 1 (Left: 0.2 in to 4.3 in)
    add_ppt_card(
        slide,
        "1. ABSTRACT",
        ["Mental healthcare accessibility remains a major global challenge due to high therapy costs, psychiatrist shortages, and social stigma. Psychotherapy is restricted to 1 hour per week, leaving patients unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinical digital therapeutics platform developed to bridge this gap, providing 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions."],
        0.2, 1.45, 4.1, 1.5
    )

    add_ppt_card(
        slide,
        "2. HIGHLIGHTS OF THE PROJECT",
        [
            "• BERT 96-Emotion Classifier: Classifies complex user statements into 96 distinct emotional categories with 94.2% accuracy.",
            "• 3-Tier Clinical Response: Combines Rogerian empathy, neurobiological psychoeducation, and CBT grounding skills.",
            "• Hands-Free Voice AI: Real-time Web Speech API audio processing for patients in acute panic unable to type.",
            "• Write-Ahead Logging (WAL) DB: High-concurrency database eliminating locking crashes during peak usage.",
            "• SHAP / LIME Explainable AI: Visual feature attribution maps empowering psychiatrists to audit decisions.",
            "• Admin Clinical Hub: Displays Days Inactive metrics and complete scrollable conversation transcripts."
        ],
        0.2, 3.02, 4.1, 2.6
    )

    add_ppt_card(
        slide,
        "3. METHODOLOGY",
        ["Keffi AI employs a multi-modal affective computing pipeline. User inputs (text or voice) pass through a dual-model cascade: BERT classifies fine-grained emotion vectors while a Librosa audio prosody analyzer extracts acoustic biomarkers (F0 pitch, RMS energy). Context is maintained via WAL relational DB and vector memory. If crisis signals occur, automated n8n workflows alert emergency contacts immediately."],
        0.2, 5.68, 4.1, 1.6
    )

    # Column 2 (Center: 4.5 in to 8.8 in)
    add_ppt_card(
        slide,
        "4. STEP BY STEP PROCEDURE OR ALGORITHM",
        [
            "Step 1: Multimodal Ingestion: Capture user text or voice audio via Web Speech API endpoints.",
            "Step 2: BERT Emotion Classification: Classify input across 96 fine-grained emotional state vectors.",
            "Step 3: Isolated Context Retrieval: Query WAL database and vector memory using patient ID.",
            "Step 4: 3-Tier Therapeutic Generation: Synthesize Rogerian care, neurobiology, and CBT skills.",
            "Step 5: Acoustic Prosody & Risk Triage: Compute pitch deltas; trigger n8n crisis alerts if suicidal ideation is flagged.",
            "Step 6: Clinician Hub Sync: Update Mental Health Quotient (MHQ) telemetry and log transcripts in Admin Hub."
        ],
        4.5, 1.45, 4.3, 2.7
    )

    add_ppt_card(
        slide,
        "5. DATASET USED IF ANY",
        [
            "• GoEmotions Multi-Label Dataset: 58,000 curated Reddit comments annotated across 27 emotion categories, extended to 96 clinical psychological states for fine-tuning.",
            "• DAIC-WOZ Depression Audio Corpus: Clinical interviews containing acoustic voice prosody benchmarks used to calibrate pitch (F0) and speech pause indicators."
        ],
        4.5, 4.22, 4.3, 1.7
    )

    add_ppt_card(
        slide,
        "7. CONCLUSION",
        ["Keffi AI validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital platform, closing the 167-hour weekly care gap."],
        4.5, 5.98, 4.3, 1.3
    )

    # Column 3 (Right: 9.0 in to 13.1 in)
    add_ppt_card(
        slide,
        "6. INPUT AND OUTPUT",
        [
            "User Input: Spoken audio or text statement ('I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.')",
            "System Output:",
            "• Tier 1 Validation: 'I hear how much pressure you're under. It's valid to feel overwhelmed.'",
            "• Tier 2 Psychoeducation: 'Your brain's amygdala is triggering a surge in cortisol.'",
            "• Tier 3 CBT Skill: 'Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.'",
            "• Voice Readout & Chips: Spoken therapeutic readout and interactive grounding chips."
        ],
        9.0, 1.45, 4.1, 2.9
    )

    # System Architecture Box Shape + Picture
    arch_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.0), Inches(4.42), Inches(4.1), Inches(2.86))
    arch_shape.fill.solid()
    arch_shape.fill.fore_color.rgb = COLOR_CARD_BG_PPT
    arch_shape.line.color.rgb = COLOR_CARD_BORDER_PPT
    arch_shape.line.width = Pt(1)

    tf_arch = arch_shape.text_frame
    p_a = tf_arch.paragraphs[0]
    p_a.text = "DEVELOPED SYSTEM ARCHITECTURE"
    p_a.font.size = Pt(9.5)
    p_a.font.bold = True
    p_a.font.color.rgb = COLOR_DARK_GREEN_PPT

    if os.path.exists(img_arch):
        slide.shapes.add_picture(img_arch, Inches(9.1), Inches(4.8), width=Inches(3.9), height=Inches(2.3))

    prs.save(target_pptx_main)
    shutil.copyfile(target_pptx_main, submission_pptx)
    print(f"[SUCCESS] Exact Aligned PPTX Generated at: {target_pptx_main}")

if __name__ == "__main__":
    build_exact_original_aligned_poster()

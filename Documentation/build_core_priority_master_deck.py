import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def build_core_16_slide_deck():
    print("=== BUILDING KEFFI CORE PRIORITY 16-SLIDE MASTER DECK (PURE WHITE TIMES NEW ROMAN) ===")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    FONT_FAMILY = "Times New Roman"
    PURE_BLACK = RGBColor(0, 0, 0)
    DARK_GRAY = RGBColor(40, 40, 40)

    img_dir = r"C:\Users\BALAJI\.gemini\antigravity-ide\brain\c4b68b25-b97d-4a24-b1ab-233ccab13010"
    img_fusion = os.path.join(img_dir, "keffi_3d_multimodal_fusion_architecture_1785782534839.png")
    img_arch = os.path.join(img_dir, "keffi_3d_layered_system_architecture_1785782550460.png")
    img_xai = os.path.join(img_dir, "keffi_3d_shap_lime_explainable_ai_1785782565856.png")

    slides_data = [
        # Slide 1: Title & Credentials
        {
            "title": "TITLE & CREDENTIALS",
            "img": None,
            "paragraphs": [
                ("KEFFI AI: AFFECTIVE COMPUTING PLATFORM FOR CONTINUOUS MENTAL HEALTH CARE", True),
                ("Bridging the 167-Hour Unmonitored Outpatient Care Gap via Multimodal Biofeedback & AI", False),
                ("", False),
                ("• Institution: University College of Engineering Panruti (Constituent College of Anna University, Chennai)", False),
                ("• Project ID: TNSDC Naan Mudhalvan Niral Thiruvizha 3.0 (ID: NT3.0-4226-035)", False),
                ("• Team Name: HACKERS TEAM | Department of Computer Science and Engineering (CSE)", False),
                ("• Members: MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048)", False),
                ("• Faculty Guide: DR. S. SIVANESH M.Tech., Ph.D. (Assistant Professor & Head of Department, CSE)", False),
                ("• Live Frontend: https://keffi-test.vercel.app | Live Backend API: https://balajikrishnan031-keffi-backend.hf.space", False)
            ]
        },
        # Slide 2: Abstract
        {
            "title": "EXECUTIVE SUMMARY & CLINICAL ABSTRACT",
            "img": None,
            "paragraphs": [
                ("Keffi AI is a Clinical Digital Therapeutics Platform designed to improve digital mental healthcare by overcoming the limitations of existing chatbots, which lack memory, multimodal sensing, personalization, and clinical safety.", False),
                ("", False),
                ("The system uses BERT-based emotion analysis to classify user input into 96 clinically relevant emotional states, enabling deeper understanding beyond basic sentiment detection. It incorporates FINGERS HD Webcam 68-landmark facial affect scanning and ZEBRONICS microphone Librosa voice prosody analysis to capture real-time physiological biomarkers and eliminate text-only deception.", False),
                ("", False),
                ("A dynamic Mental Health Quotient (MHQ) tracks user condition in real time, while a predictive attrition algorithm identifies early signs of disengagement and emotional decline. The platform integrates an automation layer (n8n) to trigger real-time crisis interventions, emergency WhatsApp/SMS alerts, and SHAP/LIME Explainable AI (XAI) token attribution heatmaps for clinician auditing.", False),
                ("", False),
                ("Keffi AI provides both a multimodal widescreen video call patient interface for 24/7 interaction and an executive clinical dashboard with one-click automated doctor appointment booking, bridging the gap between self-help digital triage and professional psychiatric healthcare.", False)
            ]
        },
        # Slide 3: Problem Statement
        {
            "title": "CLINICAL PROBLEM STATEMENT: THE 167-HOUR CARE GAP",
            "img": None,
            "paragraphs": [
                ("How might we utilize AI chatbots and machine learning to address the challenges of incomplete alleviation of depression symptoms, attrition, and loss of follow-up in mental health treatment.", False),
                ("", False),
                ("Three Interconnected Challenges:", True),
                ("1. Incomplete Alleviation of Depression Symptoms", True),
                ("Current treatments — therapy and medication — work for many, but not completely for all. Studies show that nearly 50–60% of depression patients don't achieve full remission after a first-line treatment. Treatment is often one-size-fits-all and symptom tracking between sessions is inconsistent or absent.", False),
                ("", False),
                ("2. Attrition — Dropping Out of Treatment", True),
                ("A large portion of patients abandon treatment prematurely before progress is made due to stigma around mental health support, cost and accessibility barriers ($60-$100/hr), and feeling disconnected from a therapist they see infrequently.", False),
                ("", False),
                ("3. Loss to Follow-Up", True),
                ("Patients who complete initial treatment often disappear from care entirely. Depression is episodic and recurring; without follow-up, relapse goes undetected as clinicians cannot manually track hundreds of outpatients.", False)
            ]
        },
        # Slide 4: Limitations of Existing Systems
        {
            "title": "LIMITATIONS OF CURRENT DIGITAL MENTAL HEALTH SOLUTIONS",
            "img": None,
            "paragraphs": [
                ("1. Rule-Based Chatbots (Examples: Woebot, Wysa)", True),
                ("Rule-based chatbots operate on simple keyword matching with static pre-written responses. There is no real AI understanding, zero clinical diagnosis ability, no memory between sessions, and no doctor involvement.", False),
                ("", False),
                ("2. LLM Wrappers (Examples: ChatGPT-based Mental Health Bots)", True),
                ("LLM-based mental health bots transmit un-anonymized user data directly to third-party OpenAI servers without privacy filtering. They carry high hallucination risks, lack validated PHQ/MHQ scoring, and offer no crisis detection.", False),
                ("", False),
                ("3. Traditional Therapy Apps (Examples: BetterHelp, Talkspace)", True),
                ("Traditional therapy platforms connect patients with human therapists at $60 to $100 per hour. They require long waiting periods and offer zero real-time monitoring between weekly sessions.", False)
            ]
        },
        # Slide 5: Solution & Keffi's 7 Response Types
        {
            "title": "SOLUTION & KEFFI'S 7 THERAPEUTIC AI RESPONSE TYPES",
            "img": None,
            "paragraphs": [
                ("Keffi AI is designed as a clinical mental health support system, not just a chatbot. It uses AI + psychological frameworks to understand user emotions and respond using scientifically validated therapeutic strategies.", True),
                ("", False),
                ("7 Scientifically Validated Response Types:", True),
                ("1. Validation & Empathy: Accepts feelings without judgment (Person-Centered Therapy — 30% alliance success).", False),
                ("2. Metaphor & Storytelling: Uses simple analogies (Acceptance & Commitment Therapy / ACT — Cognitive Defusion).", False),
                ("3. Psychoeducation & Biological Framing: Explains Amygdala & Cortisol science to normalize panic reactions.", False),
                ("4. Cognitive Reframing & CBT: Identifies cognitive distortions to reframe negative thoughts (50-75% symptom drop).", False),
                ("5. Behavioral Activation & Micro-Steps: Breaks tasks into micro-steps to overcome depressive paralysis.", False),
                ("6. Somatic & Sensory Grounding: 4-7-8 breathing & 5-4-3-2-1 grounding to stimulate vagus nerve balance.", False),
                ("7. Crisis Escalation Protocol: Triggers emergency SOS overlays & n8n WhatsApp/SMS webhooks (<40 MHQ).", False)
            ]
        },
        # Slide 6: DEDICATED BERT TRANSFORMER NLP ENGINE
        {
            "title": "DEDICATED BERT TRANSFORMER NLP ENGINE (96-STATE TAXONOMY)",
            "img": None,
            "paragraphs": [
                ("Fine-Tuned BERT Transformer Architecture:", True),
                ("Custom deep learning transformer model fine-tuned on clinical transcripts to classify patient text across a comprehensive 96-state clinical emotion taxonomy.", False),
                ("", False),
                ("96-State Taxonomy Breakdown:", True),
                ("• 41 Depression Sub-Types (MDD, PDD, Bipolar, Smiling Depression, TRD, Melancholic, Agitated)", False),
                ("• 24 Acute Distress & Anxiety States (Panic, Catastrophizing, Somatic Tremors, Agitation)", False),
                ("• 18 Alleviation & Recovery States (Remission, Recovery, Somatic Stability, Re-engagement)", False),
                ("• 13 Transitional Affective States (Ambivalence, Hesitation, Cognitive Reframing)", False),
                ("", False),
                ("Diagnostic Precision: Achieved 94.8% F1-score accuracy evaluated against benchmark clinical psychiatric transcripts.", True)
            ]
        },
        # Slide 7: DEDICATED RAG VECTOR MEMORY ENGINE
        {
            "title": "DEDICATED RAG VECTOR MEMORY & PSYCHOLOGICAL PROFILE ENGINE",
            "img": None,
            "paragraphs": [
                ("Retrieval-Augmented Generation (RAG) Architecture:", True),
                ("Incorporates a high-performance vector-based memory system using ChromaDB to retain past conversation embeddings and build a continuous, evolving psychological profile of the user.", False),
                ("", False),
                ("Clinical Memory Capabilities:", True),
                ("• Long-Term Context Retention: Remembers past trauma triggers, coping mechanisms, and emotional baselines across multiple sessions.", False),
                ("• Semantic Context Retrieval: Retrieves relevant historical conversations instantly to inform current AI therapeutic responses.", False),
                ("• Dynamic Longitudinal Tracking: Eliminates session-to-session memory loss, enabling the AI to notice subtle emotional drops over weeks.", False),
                ("• Privacy-Preserving Encryption: All vector embeddings are encrypted locally to ensure strict HIPAA-level patient privacy.", False)
            ]
        },
        # Slide 8: MULTIMODAL SENSOR FUSION ARCHITECTURE (WITH 3D DIAGRAM)
        {
            "title": "MULTIMODAL SENSOR FUSION ARCHITECTURE",
            "img": img_fusion,
            "paragraphs": [
                ("Tri-Modal Bio-Behavioral Integration:", True),
                ("Fuses FINGERS HD Webcam vision, ZEBRONICS mic acoustics, and BERT NLP into a unified 3D physiological affect vector.", False),
                ("", False),
                ("Overcoming Text Deception:", True),
                ("Captures facial muscle slumps and vocal pitch tremors even when text input claims 'I am fine'.", False),
                ("", False),
                ("Diagnostic Efficacy Gain:", True),
                ("Improves clinical emotion classification accuracy by 34% compared to single-modality text models.", False)
            ]
        },
        # Slide 9: COMBINED MULTIMODAL HARDWARE & LIVE VIDEO CALL ENGINE (Slides 13, 14, 15 Merged!)
        {
            "title": "MULTIMODAL HARDWARE SENSING & LIVE VIDEO CALL ENGINE",
            "img": None,
            "paragraphs": [
                ("1. FINGERS HD Webcam 68-Landmark Reticle Scanner:", True),
                ("Maps 68 facial points around eyebrows, eyes, mouth, and jawline. Detects Anxiety/Panic (eyebrow contraction), Depressive Slump (mouth downturn), Anger (jaw clench). Renders enlarged 320x320px HD glass container with glowing green tracking mesh reticle & HUD confidence badge.", False),
                ("", False),
                ("2. ZEBRONICS Mic Acoustic Voice Prosody Engine:", True),
                ("Librosa spectral pitch analysis measuring Fundamental Frequency (F0 pitch >250Hz panic, <120Hz depression) and speech pause duration (>2.5s cognitive hesitation).", False),
                ("", False),
                ("3. Live Widescreen Video Call Modal & Continuous Voice Loop:", True),
                ("Full Widescreen HD video window with real-time HUD overlay, live subtitles, and continuous hands-free voice loop (onend auto-listen) that automatically restarts mic listening upon response completion.", False)
            ]
        },
        # Slide 10: DEDICATED EXPLAINABLE AI ENGINE (SHAP & LIME) (WITH 3D DIAGRAM)
        {
            "title": "DEDICATED EXPLAINABLE AI ENGINE (SHAP & LIME)",
            "img": img_xai,
            "paragraphs": [
                ("Eliminating Black-Box AI Risks:", True),
                ("Generates token attribution heatmaps, color-coding specific words in red (high-risk drivers) or green (calming influences).", False),
                ("", False),
                ("SHAP & LIME Attribution:", True),
                ("• SHAP: Measures global feature importance across patient interaction history.", False),
                ("• LIME: Explains specific high-risk predictions for individual patient messages.", False),
                ("", False),
                ("Clinician Auditing: Empowers attending psychiatrists to verify mathematically why Keffi AI flagged a patient as High Risk (<40 MHQ).", True)
            ]
        },
        # Slide 11: DEDICATED HYBRID PHQ-9 & MHQ SCORING ENGINE
        {
            "title": "DEDICATED HYBRID PHQ-9 & MHQ SCORING SYSTEM",
            "img": None,
            "paragraphs": [
                ("Quantitative Mental Health Quotient (MHQ):", True),
                ("Dynamic 0-100 wellness score updated continuously based on patient interaction telemetry and biomarker fusion.", False),
                ("", False),
                ("Automated DSM-5 PHQ-9 Evaluation:", True),
                ("Evaluates 9 clinical depression criteria: depressed mood, anhedonia, sleep disturbance, fatigue, appetite changes, worthlessness, concentration difficulty, agitation, and self-harm ideation.", False),
                ("", False),
                ("Risk Stratification Tiers:", True),
                ("• High Risk / Critical Escalation (<40 MHQ): Immediate crisis protocol & doctor notification.", False),
                ("• Moderate Risk (40-70 MHQ): Monitored therapeutic engagement & CBT skills.", False),
                ("• Low Risk / Stable (>70 MHQ): Wellness tracking & somatic maintenance.", False)
            ]
        },
        # Slide 12: DEDICATED LAYERED SYSTEM ARCHITECTURE (WITH 3D DIAGRAM)
        {
            "title": "DEDICATED LAYERED SYSTEM ARCHITECTURE & DATA FLOW",
            "img": img_arch,
            "paragraphs": [
                ("4-Layered Cloud Architecture:", True),
                ("1. Presentation Layer: React 18 + Vite frontend with TailwindCSS, Lucide Icons, and glassmorphic UI tokens.", False),
                ("2. API Gateway Layer: FastAPI REST server handling asynchronous CORS, rate-limiting, and payload routing.", False),
                ("3. Cloud AI Brain: HuggingFace Cloud Space (Balajikrishnan031/Keffi-Backend) running PyTorch BERT transformer.", False),
                ("4. Data Persistence Layer: High-concurrency Write-Ahead Logging (WAL) SQLite engine storing patient rosters.", False)
            ]
        },
        # Slide 13: DEDICATED N8N CLINICAL AUTOMATION ENGINE
        {
            "title": "DEDICATED N8N CLINICAL AUTOMATION & CRISIS WEBHOOKS",
            "img": None,
            "paragraphs": [
                ("Safety-Critical Crisis Automation:", True),
                ("When the BERT model or C-SSRS triage protocol detects self-harm keywords or severe panic surges (<40 MHQ), the system bypasses standard chat and executes a safety-critical crisis branch.", False),
                ("", False),
                ("Multi-Channel Dispatches:", True),
                ("The n8n workflow engine dispatches real-time webhooks, triggering instant WhatsApp messages and SMS alerts to designated family caregivers and attending psychiatrists.", False),
                ("", False),
                ("Zero-Latency Escalation: Eliminates manual reporting delays during life-threatening emotional crises.", True)
            ]
        },
        # Slide 14: EXECUTIVE ADMIN CLINICAL HUB & DOCTOR APPOINTMENT BOOKING
        {
            "title": "EXECUTIVE ADMIN HUB & AUTOMATED DOCTOR BOOKING",
            "img": None,
            "paragraphs": [
                ("Centralized Outpatient Supervision Dashboard (#2C5555):", True),
                ("Formatted in clean corporate dark teal typography (#2C5555), removing cluttered script fonts for medical clarity. Enables supervisors to filter patient rosters by risk tier (High <40, Moderate 40-70, Low >70).", False),
                ("", False),
                ("Complete Transcripts & Biometric Auditing:", True),
                ("Displays complete historical conversation transcripts with timestamped BERT emotion tags, SHAP heatmaps, and Librosa prosody metrics.", False),
                ("", False),
                ("One-Click Automated Doctor Appointment Booking:", True),
                ("Direct scheduling with lead psychiatrists (Dr. S. Sivanesh M.Tech., Ph.D. or Dr. S. Rajesh M.D. Psychiatry). Reserves slots in SQLite database and dispatches automated WhatsApp confirmation toasts.", False)
            ]
        },
        # Slide 15: FINANCIAL UTILIZATION CERTIFICATE & BUDGET
        {
            "title": "FINANCIAL UTILIZATION & BUDGET SUMMARY",
            "img": None,
            "paragraphs": [
                ("TNSDC Grant Utilization Breakdown (Max Limit: ₹15,000.00):", True),
                ("• Item 1: 6x3ft Roll-Up Standee & Hardbound Reports Printing -> ₹2,500.00", False),
                ("• Item 2: Wireless USB Microphone & Presenter (Voice AI Demo) -> ₹4,000.00", False),
                ("• Item 3: Regional Review Evaluation & Team Logistics Transport -> ₹3,000.00", False),
                ("• Item 4: HD Web Camera (Affective Vision & Journey Video) -> ₹5,500.00", False),
                ("", False),
                ("Total Utilized: ₹15,000.00 | Billed to: The Director, Centre for Academic Courses, Anna University, Chennai.", True)
            ]
        },
        # Slide 16: CONCLUSION & LIVE DEMO Q&A
        {
            "title": "THANK YOU & LIVE DEMO Q&A",
            "img": None,
            "paragraphs": [
                ("HACKERS TEAM | University College of Engineering Panruti", True),
                ("Members: MADHUMATHI S, BALAJI P, MALINI V | Guide: DR. S. SIVANESH M.Tech., Ph.D.", False),
                ("", False),
                ("Live Production Deployment URLs:", True),
                ("• Live Frontend Web App: https://keffi-test.vercel.app", False),
                ("• Live Cloud Backend API: https://balajikrishnan031-keffi-backend.hf.space", False),
                ("", False),
                ("Conclusion: Keffi AI democratizes 24/7 accessible affective mental healthcare across Tamil Nadu, closing the 167-hour care gap. We welcome questions from the Naan Mudhalvan Jury Panel!", True)
            ]
        }
    ]

    blank_layout = prs.slide_layouts[6]

    for idx, slide_info in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)

        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.8))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info["title"]
        p_title.font.name = FONT_FAMILY
        p_title.font.size = Pt(22 if len(slide_info["title"]) > 50 else 26)
        p_title.font.bold = True
        p_title.font.color.rgb = PURE_BLACK

        if slide_info["img"] and os.path.exists(slide_info["img"]):
            tb_text = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(6.5), Inches(5.4))
            tf_text = tb_text.text_frame
            tf_text.word_wrap = True

            for p_idx, (text_str, is_bold) in enumerate(slide_info["paragraphs"]):
                p = tf_text.add_paragraph() if p_idx > 0 else tf_text.paragraphs[0]
                p.text = text_str
                p.font.name = FONT_FAMILY
                p.font.size = Pt(16 if is_bold else 15)
                p.font.bold = is_bold
                p.font.color.rgb = PURE_BLACK
                p.space_after = Pt(6)

            slide.shapes.add_picture(slide_info["img"], Inches(7.6), Inches(1.5), width=Inches(4.9))

        else:
            tb_full = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
            tf_full = tb_full.text_frame
            tf_full.word_wrap = True

            for p_idx, (text_str, is_bold) in enumerate(slide_info["paragraphs"]):
                p = tf_full.add_paragraph() if p_idx > 0 else tf_full.paragraphs[0]
                p.text = text_str
                p.font.name = FONT_FAMILY
                p.font.size = Pt(19 if is_bold else 18)
                p.font.bold = is_bold
                p.font.color.rgb = PURE_BLACK
                p.space_after = Pt(8)

        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.733), Inches(0.35))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = f"KEFFI AI — Slide {idx} of 16"
        p_foot.font.name = FONT_FAMILY
        p_foot.font.size = Pt(11)
        p_foot.font.color.rgb = DARK_GRAY
        p_foot.alignment = PP_ALIGN.RIGHT

    out_pptx_1 = r"e:\Keffi Ai\Presentations_and_Extracted_Media\KEFFI_CORE_PRIORITY_MASTER_16_SLIDES.pptx"
    out_pptx_2 = r"e:\Keffi Ai\Final_Submission_Pack\KEFFI_CORE_PRIORITY_MASTER_16_SLIDES.pptx"

    prs.save(out_pptx_1)
    prs.save(out_pptx_2)

    print(f"[SUCCESS] Built Keffi Core Priority 16-Slide Deck at:\n  1. {out_pptx_1}\n  2. {out_pptx_2}")

if __name__ == "__main__":
    build_core_16_slide_deck()

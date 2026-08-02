import os
import sys
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_pptx_deck_and_poster():
    print("=== GENERATING MASTER PPTX FILES (KEFFI_POSTER.PPTX & KEFFI_FINAL_REVIEW_PRESENTATION.PPTX) ===")
    
    media_dir = r'e:\Keffi Ai\Presentations_and_Extracted_Media'
    pack_dir = r'e:\Keffi Ai\Final_Submission_Pack'
    os.makedirs(media_dir, exist_ok=True)
    os.makedirs(pack_dir, exist_ok=True)

    img_logo_left = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_1.png'
    img_logo_right = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_3.jpg'
    img_hero = r'e:\Keffi Ai\Documentation\extracted_report_images\shot_1_landing_hero.png'
    img_arch = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_4.jpg'

    COLOR_DARK_GREEN = RGBColor(15, 56, 56)      # #0F3838
    COLOR_TEAL_HEAD = RGBColor(13, 80, 80)       # #0D5050
    COLOR_SLATE_TEXT = RGBColor(30, 41, 59)      # #1E293B
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_LIGHT_BG = RGBColor(240, 247, 245)     # #F0F7F5
    COLOR_BORDER_GREEN = RGBColor(143, 169, 137) # #8FA989

    # =========================================================================
    # PART 1: SINGLE-SLIDE A3 POSTER PPTX (KEFFI_POSTER.pptx)
    # =========================================================================
    prs_poster = Presentation()
    # A3 Landscape Slide Size: 16.54 inches x 11.69 inches
    prs_poster.slide_width = Inches(16.54)
    prs_poster.slide_height = Inches(11.69)
    blank_layout = prs_poster.slide_layouts[6]
    slide_p = prs_poster.slides.add_slide(blank_layout)

    # 1. Header Banner Shape
    banner = slide_p.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.4), Inches(15.74), Inches(1.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_DARK_GREEN
    banner.line.color.rgb = COLOR_DARK_GREEN

    tf = banner.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "KEFFI AI – A MENTAL HEALTH CHATBOT"
    p0.font.size = Pt(24)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_WHITE
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "UNIVERSITY COLLEGE OF ENGINEERING PANRUTI"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering"
    p2.font.size = Pt(10)
    p2.font.italic = True
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "TEAM NAME: HACKERS TEAM   |   MEMBERS: MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048)   |   GUIDE: DR. S. SIVANESH M.Tech., Ph.D."
    p3.font.size = Pt(9.5)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_WHITE
    p3.alignment = PP_ALIGN.CENTER

    # Logos on banner
    if os.path.exists(img_logo_left):
        slide_p.shapes.add_picture(img_logo_left, Inches(0.6), Inches(0.5), width=Inches(1.3), height=Inches(1.3))
    if os.path.exists(img_logo_right):
        slide_p.shapes.add_picture(img_logo_right, Inches(14.6), Inches(0.5), width=Inches(1.3), height=Inches(1.3))

    # Helper function for text boxes
    def add_section_box(slide, heading, content_lines, left_in, top_in, width_in, height_in):
        tx_box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
        tf_box = tx_box.text_frame
        tf_box.word_wrap = True
        
        ph = tf_box.paragraphs[0]
        ph.text = heading.upper()
        ph.font.size = Pt(13)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_TEAL_HEAD
        ph.alignment = PP_ALIGN.LEFT
        
        for line in content_lines:
            pl = tf_box.add_paragraph()
            pl.text = line
            pl.font.size = Pt(9.5)
            pl.font.color.rgb = COLOR_SLATE_TEXT
            pl.space_after = Pt(4)
            pl.alignment = PP_ALIGN.LEFT

    # Column 1 (Left: 0.4 to 5.4 in)
    add_section_box(
        slide_p,
        "1. ABSTRACT",
        ["Mental healthcare accessibility remains a critical global challenge due to high therapy costs, psychiatrist shortages, and social stigma. Standard psychotherapy is restricted to 1 hour per week, leaving patients completely unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinical digital therapeutics platform developed to bridge this gap, providing 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions."],
        0.4, 2.1, 5.0, 2.2
    )

    add_section_box(
        slide_p,
        "2. HIGHLIGHTS OF THE PROJECT",
        [
            "• Fine-Tuned BERT 96-Emotion Classifier: Classifies complex user statements into 96 distinct emotional categories with 94.2% accuracy.",
            "• 3-Tier Clinical Therapeutic Response: Constructs responses combining Rogerian empathy, neurobiological psychoeducation, and CBT grounding skills.",
            "• Hands-Free Voice-to-Voice AI: Integrated real-time Web Speech API audio processing for patients in acute panic unable to type.",
            "• Write-Ahead Logging (WAL) DB Engine: High-concurrency database architecture eliminating locking crashes during peak usage.",
            "• Explainable AI (SHAP / LIME): Visual feature attribution maps empowering psychiatrists to audit decisions.",
            "• Admin Clinical Hub: Displays Days Inactive metrics and complete scrollable conversation transcripts."
        ],
        0.4, 4.5, 5.0, 6.5
    )

    # Column 2 (Center: 5.7 to 10.9 in)
    add_section_box(
        slide_p,
        "3. METHODOLOGY",
        ["Keffi AI employs a multi-modal affective computing pipeline. User inputs (text or voice audio) pass through a dual-model processing cascade. The BERT transformer classifies fine-grained emotion vectors, while an audio prosody analyzer extracts acoustic biomarkers (F0 pitch, RMS energy). Context is maintained by retrieving chronological session memory from a WAL-enabled relational database and vector embedding store. If crisis signals occur, automated n8n workflows alert emergency contacts immediately."],
        5.7, 2.1, 5.2, 2.2
    )

    add_section_box(
        slide_p,
        "4. STEP BY STEP PROCEDURE OR ALGORITHM",
        [
            "Step 1: Multimodal Data Ingestion: Capture user text utterance or real-time voice audio via Web Speech API.",
            "Step 2: BERT Emotion Classification: Tokenize text input and classify across 96 fine-grained emotional state vectors.",
            "Step 3: Isolated Context Retrieval: Query WAL database and vector memory using patient ID to restore chat timeline.",
            "Step 4: 3-Tier Therapeutic Generation: Synthesize Rogerian validation, neurobiology, and CBT grounding exercises.",
            "Step 5: Acoustic Prosody & Risk Triage: Compute pitch variance deltas; trigger n8n crisis alerts if suicidal ideation is flagged.",
            "Step 6: Clinician Hub Sync: Update Mental Health Quotient (MHQ) score telemetry and log transcripts in Admin Hub."
        ],
        5.7, 4.5, 5.2, 3.8
    )

    add_section_box(
        slide_p,
        "5. DATASET USED IF ANY",
        [
            "• GoEmotions Multi-Label Dataset: 58,000 curated Reddit comments annotated across 27 emotion categories, extended to 96 clinical psychological states.",
            "• DAIC-WOZ Depression Audio Corpus: Clinical interviews containing acoustic voice prosody benchmarks used to calibrate pitch (F0) and speech pause indicators."
        ],
        5.7, 8.5, 5.2, 2.5
    )

    # Column 3 (Right: 11.2 to 16.1 in)
    add_section_box(
        slide_p,
        "6. INPUT AND OUTPUT",
        [
            "User Input: Spoken voice audio or written text statement ('I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.')",
            "System Output:",
            "• Tier 1 Validation: 'I hear how much pressure you're under. It's valid to feel overwhelmed.'",
            "• Tier 2 Psychoeducation: 'Your brain's amygdala is triggering a surge in cortisol.'",
            "• Tier 3 CBT Skill: 'Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.'",
            "• Voice Readout & Chips: Spoken therapeutic readout and interactive grounding chips."
        ],
        11.2, 2.1, 4.9, 3.5
    )

    add_section_box(
        slide_p,
        "7. CONCLUSION",
        ["Keffi AI successfully validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital mental health platform, closing the 167-hour weekly care gap."],
        11.2, 5.8, 4.9, 1.8
    )

    if os.path.exists(img_hero):
        slide_p.shapes.add_picture(img_hero, Inches(11.2), Inches(7.8), width=Inches(4.9), height=Inches(2.2))
    elif os.path.exists(img_arch):
        slide_p.shapes.add_picture(img_arch, Inches(11.2), Inches(7.8), width=Inches(4.9), height=Inches(2.2))

    poster_pptx_path = os.path.join(media_dir, "KEFFI_POSTER.pptx")
    prs_poster.save(poster_pptx_path)
    shutil.copyfile(poster_pptx_path, os.path.join(pack_dir, "3_Project_Poster.pptx"))
    print(f"[SUCCESS] Single-Slide A3 Poster PPTX Created: {poster_pptx_path}")

    # =========================================================================
    # PART 2: FULL 16-SLIDE FINAL PROJECT REVIEW PRESENTATION (KEFFI_FINAL_REVIEW_PRESENTATION.pptx)
    # =========================================================================
    prs_deck = Presentation()
    prs_deck.slide_width = Inches(13.33)  # Standard Widescreen 16:9
    prs_deck.slide_height = Inches(7.5)

    def add_deck_slide(title_text, bullet_items):
        slide = prs_deck.slides.add_slide(prs_deck.slide_layouts[6])
        
        # Header banner shape
        b = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
        b.fill.solid()
        b.fill.fore_color.rgb = COLOR_DARK_GREEN
        b.line.color.rgb = COLOR_DARK_GREEN
        
        tf_b = b.text_frame
        p_b = tf_b.paragraphs[0]
        p_b.text = title_text.upper()
        p_b.font.size = Pt(20)
        p_b.font.bold = True
        p_b.font.color.rgb = COLOR_WHITE
        p_b.alignment = PP_ALIGN.LEFT
        
        # Add logo small on right corner
        if os.path.exists(img_logo_right):
            slide.shapes.add_picture(img_logo_right, Inches(12.3), Inches(0.15), width=Inches(0.8), height=Inches(0.8))

        # Main text container
        tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.73), Inches(5.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        for i, item in enumerate(bullet_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_SLATE_TEXT
            p.space_after = Pt(12)

    # 1. Title Slide
    s1 = prs_deck.slides.add_slide(prs_deck.slide_layouts[6])
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK_GREEN
    tf1 = bg1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "\nKEFFI AI – A MENTAL HEALTH CHATBOT"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER
    
    p = tf1.add_paragraph()
    p.text = "A Clinical Digital Therapeutics Platform for Continuous 24/7 Affective Support\n"
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = COLOR_BORDER_GREEN
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "UNIVERSITY COLLEGE OF ENGINEERING PANRUTI\n(A Constituent College of Anna University, Chennai) • Department of CSE\n"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    p = tf1.add_paragraph()
    p.text = "TEAM: HACKERS TEAM   |   MEMBERS: MADHUMATHI S, BALAJI P, MALINI V   |   GUIDE: DR. S. SIVANESH M.Tech., Ph.D."
    p.font.size = Pt(12)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.CENTER

    # 2. Executive Summary
    add_deck_slide("1. EXECUTIVE ABSTRACT & MISSION", [
        "• Mental healthcare accessibility remains a critical global emergency due to therapy costs, psychiatrist shortages, and social stigma.",
        "• Standard psychotherapy provides 1 hour of weekly care, leaving patients unmonitored during the remaining 167 hours of weekly vulnerability.",
        "• Keffi AI bridges this 167-hour care gap by combining fine-tuned BERT transformer emotion classification, 3-tier clinical therapeutic response, and real-time voice interaction.",
        "• The system empowers self-guided patient recovery while keeping attending psychiatrists informed through real-time risk tracking and transcript auditing."
    ])

    # 3. Problem Statement
    add_deck_slide("2. PROBLEM STATEMENT & THE 167-HOUR CARE GAP", [
        "• High Therapy Costs & Stigma: Over 70% of individuals in psychological distress receive zero clinical treatment.",
        "• The 167-Hour Care Gap: Symptoms escalate undetected between weekly counseling sessions.",
        "• Limitations of Existing Chatbots: Rule-based chatbots lack long-term memory, emotional personalization, and safety fallback mechanisms.",
        "• Lack of Clinician Transparency: Black-box AI models fail to provide explainable decision rationale for medical oversight."
    ])

    # 4. Objectives
    add_deck_slide("3. PROJECT OBJECTIVES", [
        "1. Build a Fine-Tuned BERT 96-Emotion Classifier achieving >94% accuracy for fine-grained affective state detection.",
        "2. Implement a 3-Tier Therapeutic Response Architecture (Rogerian Validation + Neurobiological Psychoeducation + CBT Skills).",
        "3. Integrate Hands-Free Voice-to-Voice AI using Web Speech API for users in acute panic.",
        "4. Develop a High-Concurrency WAL Database Engine to handle simultaneous reads and writes without locking crashes.",
        "5. Feature Explainable AI (SHAP & LIME) and an Admin Clinical Hub for complete transcript inspection and patient tracking."
    ])

    # 5. System Architecture
    add_deck_slide("4. SYSTEM ARCHITECTURE & DATA FLOW", [
        "• Multimodal Data Ingestion: Accepts text utterances and raw voice audio signals.",
        "• Dual-Model Inference Cascade: BERT fine-tuned transformer classifies 96 emotional states; Librosa extracts acoustic pitch (F0) & pause biomarkers.",
        "• Isolated Persistence Layer: WAL relational database and vector memory store session history securely tagged by patient ID.",
        "• Safety Automation: Automated n8n crisis triggers send immediate emergency SOS alerts when suicidal ideation is detected."
    ])

    # 6. BERT Model & Emotion Engine
    add_deck_slide("5. BERT TRANSFORMER EMOTION CLASSIFICATION ENGINE", [
        "• Fine-tuned over GoEmotions multi-label corpus (58,000 annotated comments).",
        "• Expanded from 27 baseline categories to 96 fine-grained clinical psychological states (e.g. panic, burnout, catastrophizing, grief).",
        "• Achieved 94.2% top-3 classification accuracy with sub-600ms latency.",
        "• Enables dynamic temperature modulation for deterministic crisis triage."
    ])

    # 7. 3-Tier Therapeutic Response Architecture
    add_deck_slide("6. 3-TIER CLINICAL THERAPEUTIC RESPONSE ARCHITECTURE", [
        "• Tier 1: Empathetic Validation (Rogerian Care) – Mirrors user emotion without judgment to establish psychological safety.",
        "• Tier 2: Psychoeducation (Neurobiological Insights) – Explains biological mechanisms (amygdala fight-or-flight, cortisol/adrenaline surges) to de-stigmatize distress.",
        "• Tier 3: CBT Micro-Intervention – Delivers immediate physical grounding skills (5-4-3-2-1 sensory mapping, 4-7-8 breathing) + 3 interactive quick-action prompt chips."
    ])

    # 8. Hands-Free Voice AI
    add_deck_slide("7. HANDS-FREE REAL-TIME VOICE-TO-VOICE AI", [
        "• Engineered for patients experiencing acute anxiety, motor tremors, or sensory overload unable to type on physical keyboards.",
        "• Continuous Web Speech API acoustic speech-to-text capture in real time.",
        "• Speech synthesis audio renderer emulating warm, empathetic voice pitch and pacing.",
        "• Bidirectional hands-free audio loop operates continuously without requiring manual button clicks."
    ])

    # 9. Voice Prosody Acoustic Tracking
    add_deck_slide("8. VOICE PROSODY ACOUSTIC TRACKING", [
        "• Detects subtle signs of depression or psychomotor retardation masked in text.",
        "• Librosa signal processing extracts Fundamental Frequency (F0 pitch contours), RMS Energy, Speech Rate, and Pause Durations.",
        "• Detects monotone pitch variance and >40% speech pause increases to flag depressive escalation.",
        "• Feeds acoustic delta vectors into the Mental Health Quotient (MHQ) scoring pipeline."
    ])

    # 10. High-Concurrency WAL Database
    add_deck_slide("9. HIGH-CONCURRENCY WAL DATABASE ENGINE", [
        "• Configured relational database in Write-Ahead Logging (WAL) mode.",
        "• Appends new writes to a separate WAL log file, allowing concurrent read queries to proceed without database locking crashes.",
        "• Tuned connection pooling and synchronous disk flush policies for ACID transactional compliance.",
        "• Ensures zero-downtime reliability under peak user concurrency."
    ])

    # 11. Explainable AI (SHAP & LIME)
    add_deck_slide("10. EXPLAINABLE ARTIFICIAL INTELLIGENCE (SHAP & LIME)", [
        "• Eliminates black-box diagnostic risks by calculating mathematical Shapley token attribution values.",
        "• Measures exact marginal contribution of individual words (e.g. 'overwhelmed', 'can't sleep') toward assigned emotion categories.",
        "• Visualizes token-level feature weight maps directly within the clinician dashboard.",
        "• Fosters clinical trust and enables attending psychiatrists to audit AI therapeutic decisions."
    ])

    # 12. Admin Clinical Hub
    add_deck_slide("11. ADMIN CLINICAL HUB & PATIENT TRACKING", [
        "• Real-Time Clinician Dashboard displaying patient rosters, distress alerts, and emotion distribution analytics.",
        "• Days Inactive Metric Counter: Tracks patient engagement gaps to flag individuals showing sudden drop-offs or emotional decline.",
        "• Admin Complete Transcript Inspection Viewer: Enables full, scrollable A-to-Z chat history auditing for selected patients.",
        "• Enables attending psychiatrists to validate AI interventions and step in manually when crisis alerts trigger."
    ])

    # 13. System Outcome & Screenshots
    add_deck_slide("12. SYSTEM OUTCOME & LIVE SCREENSHOTS", [
        "• Landing Page Hero, Silent Crisis 167-Hour Gap, and Full-Stack Tech Architecture.",
        "• Patient Login & Identity Verification Modal.",
        "• Keffi Chatting Sanctuary Room with Mic Icon & Peace Log Sidebar.",
        "• Admin Clinical Hub Doctor Dashboard Roster & Days Inactive Metrics.",
        "• Complete Scrollable Transcript Viewer."
    ])

    # 14. Performance Metrics
    add_deck_slide("13. SYSTEM PERFORMANCE & EVALUATION METRICS", [
        "• BERT Emotion Accuracy: 94.2% top-3 accuracy across 96 clinical categories.",
        "• Response Latency: Sub-600ms end-to-end API roundtrip response delivery.",
        "• System Availability: 99.9% uptime via dual-model cascade failover routing.",
        "• Concurrency Stress Test: 100+ parallel chat sessions handled without WAL database locking."
    ])

    # 15. Conclusion & Future Roadmap
    add_deck_slide("14. CONCLUSION & FUTURE ROADMAP", [
        "• Conclusion: Keffi AI successfully validates multi-modal affective computing, hands-free voice therapeutics, WAL high-concurrency storage, and clinical tracking.",
        "• Future Enhancements:",
        "  - Multi-lingual regional voice models (Tamil, Hindi).",
        "  - Hospital EHR system integration (HL7 / FHIR standards).",
        "  - Continuous PPG optical sensor analytics."
    ])

    deck_pptx_path = os.path.join(media_dir, "KEFFI_FINAL_REVIEW_PRESENTATION.pptx")
    prs_deck.save(deck_pptx_path)
    shutil.copyfile(deck_pptx_path, os.path.join(pack_dir, "KEFFI_FINAL_REVIEW_PRESENTATION.pptx"))
    print(f"[SUCCESS] 16-Slide Final Review Presentation PPTX Created: {deck_pptx_path}")

if __name__ == "__main__":
    build_pptx_deck_and_poster()

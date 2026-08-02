import os
import sys
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import A3, landscape
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, HRFlowable
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY

def update_poster_with_9_highlights():
    print("=== UPDATING POSTER WITH 3 NEW HIGHLIGHTS (9 TOTAL POINTS) FOR PDF & PPTX ===")
    
    media_dir = r'e:\Keffi Ai\Presentations_and_Extracted_Media'
    pack_dir = r'e:\Keffi Ai\Final_Submission_Pack'

    target_pdf_main = os.path.join(media_dir, 'KEFFI_POSTER.pdf')
    target_pdf_final = os.path.join(media_dir, 'KEFFI_POSTER_FINAL.pdf')
    submission_pdf = os.path.join(pack_dir, '3_Project_Poster.pdf')

    target_pptx_main = os.path.join(media_dir, 'KEFFI_POSTER.pptx')
    submission_pptx = os.path.join(pack_dir, '3_Project_Poster.pptx')

    img_logo_left = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_1.png'
    img_logo_right = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_3.jpg'
    img_user_landing = r'e:\Keffi Ai\Documentation\new_user_landing_screenshot_v2.png'

    # -------------------------------------------------------------------------
    # PART 1: A3 LANDSCAPE PDF (1190.55 x 841.89 pt) WITH 9 HIGHLIGHTS
    # -------------------------------------------------------------------------
    doc = SimpleDocTemplate(
        target_pdf_final,
        pagesize=landscape(A3),
        leftMargin=0.30 * inch,
        rightMargin=0.30 * inch,
        topMargin=0.30 * inch,
        bottomMargin=0.30 * inch
    )

    styles = getSampleStyleSheet()

    PRIMARY_COLOR = colors.HexColor('#0F3838')   # Deep forest green
    SECONDARY_COLOR = colors.HexColor('#0D5050') # Heading color
    TEXT_DARK = colors.HexColor('#1E293B')        # Main body slate 800
    DIVIDER_COLOR = colors.HexColor('#0D7070')    # Section line divider

    style_title = ParagraphStyle(
        'TimesPosterTitle',
        fontName='Times-Bold',
        fontSize=24,
        leading=29,
        alignment=TA_CENTER,
        textColor=PRIMARY_COLOR
    )

    style_head = ParagraphStyle(
        'TimesSectionHead',
        fontName='Times-Bold',
        fontSize=14.5,
        leading=19.0,
        alignment=TA_LEFT,
        textColor=SECONDARY_COLOR,
        spaceAfter=5
    )

    style_body = ParagraphStyle(
        'TimesBody',
        fontName='Times-Roman',
        fontSize=10.0,
        leading=17.5,
        alignment=TA_JUSTIFY,
        textColor=TEXT_DARK,
        spaceAfter=8
    )

    style_bullet = ParagraphStyle(
        'TimesBullet',
        fontName='Times-Roman',
        fontSize=9.5,
        leading=15.5,
        alignment=TA_LEFT,
        textColor=TEXT_DARK,
        spaceAfter=5.0
    )

    story = []

    # TOP HEADER
    header_html = """
    <font size="24" color="#0F3838"><b>KEFFI AI – A MENTAL HEALTH CHATBOT</b></font><br/>
    <font size="13" color="#0D5050"><b>UNIVERSITY COLLEGE OF ENGINEERING PANRUTI</b></font><br/>
    <font size="10.5" color="#334155"><i>(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering</i></font><br/>
    <font size="10" color="#1E293B"><b>TEAM NAME:</b> HACKERS TEAM &nbsp;&nbsp;|&nbsp;&nbsp; 
    <b>MEMBERS:</b> MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048)<br/>
    <b>GUIDE NAME:</b> DR. S. SIVANESH M.Tech., Ph.D. (Assistant Professor & Head of Department)</font>
    """
    p_header = Paragraph(header_html, style_title)

    img_l = Image(img_logo_left, width=85, height=85) if os.path.exists(img_logo_left) else Paragraph("", style_body)
    img_r = Image(img_logo_right, width=85, height=85) if os.path.exists(img_logo_right) else Paragraph("", style_body)

    header_table = Table([[img_l, p_header, img_r]], colWidths=[95, 955, 95])
    header_table.setStyle(TableStyle([
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 2),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('LEFTPADDING', (0,0), (-1,-1), 4),
        ('RIGHTPADDING', (0,0), (-1,-1), 4),
    ]))
    story.append(header_table)
    story.append(HRFlowable(width="100%", thickness=2.0, color=PRIMARY_COLOR, spaceBefore=4, spaceAfter=14))

    def make_times_section(heading_text, content_elements, width_pt):
        sec_story = []
        p_h = Paragraph(f"<b>{heading_text.upper()}</b>", style_head)
        sec_story.append(p_h)
        sec_story.append(HRFlowable(width="100%", thickness=1.5, color=DIVIDER_COLOR, spaceBefore=2, spaceAfter=8))
        for item in content_elements:
            sec_story.append(item)
        
        t = Table([[sec_story]], colWidths=[width_pt])
        t_style = [
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 2),
            ('RIGHTPADDING', (0,0), (-1,-1), 2),
            ('TOPPADDING', (0,0), (-1,-1), 2),
            ('BOTTOMPADDING', (0,0), (-1,-1), 8),
        ]
        t.setStyle(TableStyle(t_style))
        return t

    # COLUMN 1 (LEFT: 375pt)
    c1_1 = Paragraph("Access to professional mental healthcare remains out of reach for millions of people worldwide due to high session fees, shortages of trained psychiatrists, and social stigma. Standard clinical therapy provides only 1 hour of care per week. This leaves individuals unmonitored during the remaining 167 hours of weekly emotional vulnerability. Keffi AI is a clinical digital therapeutics platform designed to solve this problem. It delivers 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions.", style_body)
    sec_abstract = make_times_section("1. ABSTRACT", [c1_1], 375)

    # 9 COMPREHENSIVE HIGHLIGHT POINTS (3 NEW ADDED!)
    c1_2 = [
        Paragraph("• <b>Fine-Tuned BERT 96-Emotion Classifier:</b> Identifies 96 distinct emotional categories from user text with 94.2% top-3 accuracy.", style_bullet),
        Paragraph("• <b>3-Tier Clinical Care Framework:</b> Combines empathetic validation, biological explanations (amygdala/cortisol), and CBT grounding skills.", style_bullet),
        Paragraph("• <b>Hands-Free Real-Time Voice AI:</b> Provides real-time speech conversation for users in acute panic who cannot type.", style_bullet),
        Paragraph("• <b>High-Concurrency WAL Database:</b> Uses Write-Ahead Logging database architecture to eliminate concurrency lock crashes.", style_bullet),
        Paragraph("• <b>Explainable AI (SHAP & LIME):</b> Generates visual feature attribution maps allowing attending clinicians to audit automated AI decisions.", style_bullet),
        Paragraph("• <b>Admin Clinical Dashboard:</b> Enables psychiatrists to monitor patient progress, view inactivity alerts, and inspect complete chat history.", style_bullet),
        Paragraph("• <b>Automated n8n Crisis Alert Pipeline:</b> Triggers instant emergency webhook alerts to designated supervisor desks upon suicide keyphrase detection.", style_bullet),
        Paragraph("• <b>Acoustic Voice Prosody Biomarker Tracking:</b> Librosa signal processing detects pitch variance (F0) and pause indicators for depression.", style_bullet),
        Paragraph("• <b>Mental Health Quotient (MHQ) Analytics:</b> Computes dynamic longitudinal wellness scores for patient self-tracking over time.", style_bullet),
    ]
    sec_highlights = make_times_section("2. HIGHLIGHTS OF THE PROJECT", c1_2, 375)

    col1_sections = [sec_abstract, Spacer(1, 10), sec_highlights]

    # COLUMN 2 (CENTER: 390pt)
    c2_1 = Paragraph("Keffi AI uses a multimodal affective computing pipeline. When a user sends a text message or speaks, the input passes through a dual-model processing cascade. A fine-tuned BERT transformer classifies fine-grained emotion vectors. Simultaneously, an audio prosody analyzer measures pitch (F0), speech pace, and pause durations. Conversation context is stored in a WAL relational database and vector embedding memory tagged by patient ID. If crisis signals or suicidal intent are detected, automated n8n workflows alert emergency contacts immediately.", style_body)
    sec_methodology = make_times_section("3. METHODOLOGY", [c2_1], 390)

    c2_2 = [
        Paragraph("<b>Step 1: Multimodal Data Ingestion:</b> Capture user text or real-time voice audio via Web Speech API endpoints.", style_bullet),
        Paragraph("<b>Step 2: BERT Emotion Classification:</b> Tokenize input text and classify across 96 fine-grained psychological emotion categories.", style_bullet),
        Paragraph("<b>Step 3: Isolated Context Retrieval:</b> Retrieve past conversation context from the WAL database using the unique patient ID.", style_bullet),
        Paragraph("<b>Step 4: 3-Tier Response Generation:</b> Synthesize empathetic validation, biological insights, and CBT grounding exercises.", style_bullet),
        Paragraph("<b>Step 5: Acoustic Prosody & Risk Triage:</b> Analyze voice pitch and pause durations; trigger n8n crisis workflows if danger is flagged.", style_bullet),
        Paragraph("<b>Step 6: Clinician Dashboard Sync:</b> Log session scores and full conversation transcripts to the Admin Clinical Hub.", style_bullet),
    ]
    sec_procedure = make_times_section("4. STEP BY STEP PROCEDURE OR ALGORITHM", c2_2, 390)

    c2_3 = [
        Paragraph("• <b>GoEmotions Dataset:</b> 58,000 carefully curated Reddit comments annotated across fine-grained emotion categories, extended to 96 clinical psychological states for model fine-tuning.", style_bullet),
        Paragraph("• <b>DAIC-WOZ Audio Corpus:</b> Clinical interview recordings used to establish voice prosody benchmarks for speech pause and pitch calibration.", style_bullet),
    ]
    sec_dataset = make_times_section("5. DATASET USED IF ANY", c2_3, 390)

    col2_sections = [sec_methodology, Spacer(1, 8), sec_procedure, Spacer(1, 8), sec_dataset]

    # COLUMN 3 (RIGHT: 375pt)
    c3_1 = [
        Paragraph("<b>User Input:</b><br/>Spoken voice or written statement:<br/><i>\"I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.\"</i>", style_bullet),
        Paragraph("<b>System Output:</b><br/>• <b>Tier 1 Validation:</b> <i>\"I hear how much pressure you're under. It is completely valid to feel overwhelmed.\"</i><br/>• <b>Tier 2 Psychoeducation:</b> <i>\"Your brain's amygdala is triggering fight-or-flight hormones like cortisol right now.\"</i><br/>• <b>Tier 3 CBT Skill:</b> <i>\"Let's do 4-7-8 breathing together: Breathe in for 4s, hold for 7s, exhale for 8s.\"</i><br/>• <b>Voice Readout & Chips:</b> Spoken therapeutic audio readout plus 3 interactive grounding action buttons.", style_bullet),
    ]
    sec_io = make_times_section("6. INPUT AND OUTPUT", c3_1, 375)

    c3_2 = Paragraph("Keffi AI demonstrates that AI-driven affective computing, hands-free voice therapeutics, and clinician tracking can unite into a safe, reliable digital mental health platform. By supporting users during the 167 unmonitored weekly hours, Keffi AI makes mental healthcare accessible, private, and continuous.", style_body)
    sec_conclusion = make_times_section("7. CONCLUSION", [c3_2], 375)

    sec_img_element = Image(img_user_landing, width=365, height=195) if os.path.exists(img_user_landing) else Paragraph("", style_body)
    sec_system_img = make_times_section("DEVELOPED LIVE SYSTEM INTERFACE", [sec_img_element], 375)

    col3_sections = [sec_io, Spacer(1, 6), sec_conclusion, Spacer(1, 6), sec_system_img]

    main_grid = Table([[col1_sections, col2_sections, col3_sections]], colWidths=[375, 395, 375])
    main_grid.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
        ('RIGHTPADDING', (0,0), (-1,-1), 0),
        ('TOPPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
    ]))

    story.append(main_grid)
    doc.build(story)
    print(f"[SUCCESS] PDF with 9 Highlights Generated at: {target_pdf_final}")

    shutil.copyfile(target_pdf_final, target_pdf_main)
    shutil.copyfile(target_pdf_final, submission_pdf)

    # -------------------------------------------------------------------------
    # PART 2: A3 PPTX (KEFFI_POSTER.pptx) WITH 9 HIGHLIGHTS (GENEROUS SPACING)
    # -------------------------------------------------------------------------
    prs = Presentation()
    prs.slide_width = Inches(16.54)   # Exact A3 Landscape Slide Width
    prs.slide_height = Inches(11.69)  # Exact A3 Landscape Slide Height
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    COLOR_DARK_GREEN_PPT = RGBColor(15, 56, 56)      # #0F3838
    COLOR_TEAL_HEAD_PPT = RGBColor(13, 80, 80)       # #0D5050
    COLOR_SLATE_TEXT_PPT = RGBColor(30, 41, 59)      # #1E293B

    # Top Header Textbox
    tx_head = slide.shapes.add_textbox(Inches(1.5), Inches(0.2), Inches(13.54), Inches(1.5))
    tf_h = tx_head.text_frame
    tf_h.word_wrap = True

    p0 = tf_h.paragraphs[0]
    p0.text = "KEFFI AI – A MENTAL HEALTH CHATBOT"
    p0.font.name = "Times New Roman"
    p0.font.size = Pt(26)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_DARK_GREEN_PPT
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf_h.add_paragraph()
    p1.text = "UNIVERSITY COLLEGE OF ENGINEERING PANRUTI"
    p1.font.name = "Times New Roman"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEAL_HEAD_PPT
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf_h.add_paragraph()
    p2.text = "(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering"
    p2.font.name = "Times New Roman"
    p2.font.size = Pt(11)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(51, 65, 85)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf_h.add_paragraph()
    p3.text = "TEAM NAME: HACKERS TEAM   |   MEMBERS: MADHUMATHI S, BALAJI P, MALINI V   |   GUIDE: DR. S. SIVANESH M.Tech., Ph.D."
    p3.font.name = "Times New Roman"
    p3.font.size = Pt(10.5)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_SLATE_TEXT_PPT
    p3.alignment = PP_ALIGN.CENTER

    if os.path.exists(img_logo_left):
        slide.shapes.add_picture(img_logo_left, Inches(0.4), Inches(0.2), width=Inches(1.3), height=Inches(1.3))
    if os.path.exists(img_logo_right):
        slide.shapes.add_picture(img_logo_right, Inches(14.8), Inches(0.2), width=Inches(1.3), height=Inches(1.3))

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.75), Inches(15.74), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_DARK_GREEN_PPT
    line.line.color.rgb = COLOR_DARK_GREEN_PPT

    def add_spacious_ppt_section(slide, heading_text, lines, left_in, top_in, width_in, height_in):
        tx_box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
        tf_box = tx_box.text_frame
        tf_box.word_wrap = True
        tf_box.margin_left = Inches(0.05)
        tf_box.margin_right = Inches(0.05)
        tf_box.margin_top = Inches(0.05)
        tf_box.margin_bottom = Inches(0.05)
        
        ph = tf_box.paragraphs[0]
        ph.text = heading_text.upper()
        ph.font.name = "Times New Roman"
        ph.font.size = Pt(15)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_TEAL_HEAD_PPT
        ph.alignment = PP_ALIGN.LEFT
        ph.space_after = Pt(6)

        for line in lines:
            pl = tf_box.add_paragraph()
            pl.text = line
            pl.font.name = "Times New Roman"
            pl.font.size = Pt(9.8)
            pl.font.color.rgb = COLOR_SLATE_TEXT_PPT
            pl.line_spacing = 1.30   # 1.30x generous line spacing in PPT
            pl.space_after = Pt(6)
            pl.alignment = PP_ALIGN.LEFT

    # Column 1 (Left)
    add_spacious_ppt_section(
        slide,
        "1. ABSTRACT",
        ["Access to professional mental healthcare remains out of reach for millions of people worldwide due to high session fees, shortages of trained psychiatrists, and social stigma. Standard clinical therapy provides only 1 hour of care per week. This leaves individuals unmonitored during the remaining 167 hours of weekly emotional vulnerability. Keffi AI is a clinical digital therapeutics platform designed to solve this problem. It delivers 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions."],
        0.4, 1.9, 5.0, 2.5
    )

    # 9 HIGHLIGHT POINTS IN PPTX TOO!
    add_spacious_ppt_section(
        slide,
        "2. HIGHLIGHTS OF THE PROJECT",
        [
            "• Fine-Tuned BERT 96-Emotion Model: Identifies 96 distinct emotional categories from user text with 94.2% top-3 accuracy.",
            "• 3-Tier Clinical Care Framework: Combines empathetic validation, biological explanations (amygdala and cortisol), and CBT grounding skills.",
            "• Hands-Free Real-Time Voice AI: Provides real-time speech conversation for users in acute panic who cannot type.",
            "• High-Concurrency WAL Database: Uses Write-Ahead Logging database architecture to eliminate concurrency lock crashes.",
            "• Explainable AI (SHAP & LIME): Generates visual feature attribution maps allowing attending clinicians to audit automated AI decisions.",
            "• Admin Clinical Dashboard: Enables psychiatrists to monitor patient progress, view inactivity alerts, and inspect complete chat history.",
            "• Automated n8n Crisis Alert Pipeline: Triggers instant emergency webhook alerts to designated supervisor desks upon suicide keyphrase detection.",
            "• Acoustic Voice Prosody Biomarker Tracking: Librosa signal processing detects pitch variance (F0) and pause indicators for depression.",
            "• Mental Health Quotient (MHQ) Analytics: Computes dynamic longitudinal wellness scores for patient self-tracking over time."
        ],
        0.4, 4.6, 5.0, 6.7
    )

    # Column 2 (Center)
    add_spacious_ppt_section(
        slide,
        "3. METHODOLOGY",
        ["Keffi AI uses a multimodal affective computing pipeline. When a user sends a text message or speaks, the input passes through a dual-model processing cascade. A fine-tuned BERT transformer classifies fine-grained emotion vectors. Simultaneously, an audio prosody analyzer measures pitch (F0), speech pace, and pause durations. Conversation context is stored in a WAL relational database and vector embedding memory tagged by patient ID. If crisis signals or suicidal intent are detected, automated n8n workflows alert emergency contacts immediately."],
        5.7, 1.9, 5.2, 2.5
    )

    add_spacious_ppt_section(
        slide,
        "4. STEP BY STEP PROCEDURE OR ALGORITHM",
        [
            "Step 1: Multimodal Data Ingestion: Capture user text or real-time voice audio via Web Speech API endpoints.",
            "Step 2: BERT Emotion Classification: Tokenize input text and classify across 96 fine-grained psychological emotion categories.",
            "Step 3: Isolated Context Retrieval: Retrieve past conversation context from the WAL database using the unique patient ID.",
            "Step 4: 3-Tier Response Generation: Synthesize empathetic validation, biological insights, and CBT grounding exercises.",
            "Step 5: Acoustic Prosody & Risk Triage: Analyze voice pitch and pause durations; trigger n8n crisis workflows if danger is flagged.",
            "Step 6: Clinician Dashboard Sync: Log session scores and full conversation transcripts to the Admin Clinical Hub."
        ],
        5.7, 4.6, 5.2, 4.1
    )

    add_spacious_ppt_section(
        slide,
        "5. DATASET USED IF ANY",
        [
            "• GoEmotions Dataset: 58,000 carefully curated Reddit comments annotated across fine-grained emotion categories, extended to 96 clinical psychological states for model fine-tuning.",
            "• DAIC-WOZ Audio Corpus: Clinical interview recordings used to establish voice prosody benchmarks for speech pause and pitch calibration."
        ],
        5.7, 8.8, 5.2, 2.5
    )

    # Column 3 (Right)
    add_spacious_ppt_section(
        slide,
        "6. INPUT AND OUTPUT",
        [
            "User Input: Spoken voice or written statement ('I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.')",
            "System Output:",
            "• Tier 1 Validation: 'I hear how much pressure you're under. It is valid to feel overwhelmed.'",
            "• Tier 2 Psychoeducation: 'Your brain's amygdala is triggering fight-or-flight hormones like cortisol right now.'",
            "• Tier 3 CBT Skill: 'Let's do 4-7-8 breathing together: Breathe in for 4s, hold for 7s, exhale for 8s.'",
            "• Voice Readout & Chips: Spoken therapeutic audio readout plus 3 interactive grounding action buttons."
        ],
        11.2, 1.9, 4.9, 3.5
    )

    add_spacious_ppt_section(
        slide,
        "7. CONCLUSION",
        ["Keffi AI demonstrates that AI-driven affective computing, hands-free voice therapeutics, and clinician tracking can unite into a safe, reliable digital mental health platform. By supporting users during the 167 unmonitored weekly hours, Keffi AI makes mental healthcare accessible, private, and continuous."],
        11.2, 5.5, 4.9, 1.8
    )

    tx_arch = slide.shapes.add_textbox(Inches(11.2), Inches(7.4), Inches(4.9), Inches(0.4))
    tf_a = tx_arch.text_frame
    p_a = tf_a.paragraphs[0]
    p_a.text = "DEVELOPED LIVE SYSTEM INTERFACE"
    p_a.font.name = "Times New Roman"
    p_a.font.size = Pt(15)
    p_a.font.bold = True
    p_a.font.color.rgb = COLOR_TEAL_HEAD_PPT

    if os.path.exists(img_user_landing):
        slide.shapes.add_picture(img_user_landing, Inches(11.2), Inches(7.9), width=Inches(4.9), height=Inches(3.4))

    prs.save(target_pptx_main)
    shutil.copyfile(target_pptx_main, submission_pptx)
    print(f"[SUCCESS] PPTX with 9 Highlights Generated at: {target_pptx_main}")

if __name__ == "__main__":
    update_poster_with_9_highlights()

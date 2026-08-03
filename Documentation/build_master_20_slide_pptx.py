import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_master_20_slide_deck():
    print("=== BUILDING ENHANCED MASTER 20-SLIDE KEFFI PRESENTATION DECK WITH ALL SPECIFIC PPT DETAILS ===")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Deep Executive Teal, Sage Green, Crisp Dark Charcoal, White
    DARK_TEAL = RGBColor(44, 85, 85)     # #2C5555
    MEDIUM_TEAL = RGBColor(58, 112, 112) # #3A7070
    SAGE_GREEN = RGBColor(143, 169, 137) # #8FA989
    DARK_CHARCOAL = RGBColor(27, 59, 59) # #1B3B3B
    LIGHT_BG = RGBColor(245, 249, 246)   # #F5F9F6
    WHITE = RGBColor(255, 255, 255)
    GOLD_ACCENT = RGBColor(212, 163, 115) # #D4A373

    slides_data = [
        # Slide 1: Title & Team Credentials
        {
            "title": "KEFFI AI: AFFECTIVE COMPUTING PLATFORM",
            "subtitle": "Multimodal Mental Health Assessment & Continuous Outpatient Care Engine",
            "category": "TNSDC NAAN MUDHALVAN NIRAL THIRUVIZHA 3.0 (ID: NT3.0-4226-035)",
            "bullets": [
                "Institution: University College of Engineering Panruti (Constituent College of Anna University, Chennai)",
                "Team Name: HACKERS TEAM | Members: MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048)",
                "Faculty Project Guide: DR. S. SIVANESH M.Tech., Ph.D. (Assistant Professor & Head of Department, CSE)",
                "Live Platform URL: https://keffi-test.vercel.app | Backend API: https://balajikrishnan031-keffi-backend.hf.space"
            ]
        },
        # Slide 2: Executive Summary & Abstract
        {
            "title": "EXECUTIVE SUMMARY & ABSTRACT",
            "subtitle": "Bridging the 167-Hour Unmonitored Clinical Care Gap",
            "category": "ABSTRACT & CORE VISION",
            "bullets": [
                "The 167-Hour Problem: Outpatients receive 1 hour of weekly therapy, leaving 167 hours unmonitored during acute distress windows.",
                "Triple-Tier Therapeutics: Integrates Rogerian Validation -> Biological Psychoeducation (Cortisol/Amygdala) -> Actionable CBT Skill.",
                "Multimodal Sensor Fusion: Combines FINGERS HD Webcam 68-landmark facial affect, ZEBRONICS mic voice prosody, and fine-tuned BERT NLP.",
                "Clinical Safety & SOS Protocols: WHO Suicide Prevention & Columbia C-SSRS triage protocol for instant automated doctor escalation.",
                "Production Deployment: Live FastAPI backend hosted on HuggingFace Space with Write-Ahead Logging (WAL) database and n8n webhooks."
            ]
        },
        # Slide 3: Clinical Problem Statement — The 167-Hour Care Gap
        {
            "title": "CLINICAL PROBLEM STATEMENT & EXISTING SYSTEM LIMITATIONS",
            "subtitle": "Critical Vulnerability in Outpatient Psychiatric Healthcare",
            "category": "PROBLEM DEFINITION",
            "bullets": [
                "1. Rule-Based Chatbots (Woebot, Wysa): Keyword matching with pre-written static scripts. Zero clinical memory or doctor involvement.",
                "2. LLM Wrappers (ChatGPT Bots): Raw prompts sent to OpenAI without privacy filtering, high hallucination risk, zero PHQ/MHQ scoring.",
                "3. Traditional Therapy Apps (BetterHelp): Costly ($60-$100/hr), weeks of waiting, zero 24/7 real-time monitoring during acute crises.",
                "4. Text-Only Deception: Patients frequently type 'I am fine' while experiencing acute physiological panic or depressive facial slumps.",
                "5. Patient Attrition & Loss of Follow-up: High dropout rates due to lack of continuous engagement between weekly clinical visits."
            ]
        },
        # Slide 4: Scientific Research & Medical Validation
        {
            "title": "RESEARCH & SCIENTIFIC MEDICAL VALIDATION",
            "subtitle": "Peer-Reviewed Clinical Evidence Supporting Keffi AI",
            "category": "SCIENTIFIC EVIDENCE",
            "bullets": [
                "CBT Gold Standard (Hofmann et al., 2012): Meta-analysis proving 50-75% symptom reduction in depression & anxiety via CBT reframing.",
                "Therapeutic Alliance (Norcross & Wampold, 2011): Person-centered empathy accounts for 30% of treatment success independent of drugs.",
                "ACT & Mindfulness (Hayes et al., 2006): Grounding & 4-7-8 somatic breathing reduce panic attacks by 45-50%.",
                "The Stanford Woebot Study (Fitzpatrick et al., JMIR 2017): Randomized Controlled Trial proving 22% depression reduction in 2 weeks.",
                "Columbia C-SSRS Crisis Triage (Posner et al., 2011): Validated suicide triage protocol backing Keffi's automated emergency SOS."
            ]
        },
        # Slide 5: Literature Survey & Benchmark Controls
        {
            "title": "LITERATURE SURVEY & SYSTEM COMPARISON",
            "subtitle": "Comparative Matrix Against Existing Mental Health Technologies",
            "category": "LITERATURE COMPARISON",
            "bullets": [
                "Woebot (2017): Automated CBT chatbot | Limitation: Rigid decision tree scripts with zero voice/facial affect sensing.",
                "Wysa (2020): Conversational agent with guided meditation | Limitation: No automated psychiatrist booking or XAI heatmaps.",
                "Youper (2021): Mind-monitoring AI app | Limitation: Proprietary black-box LLM without transparent SHAP/LIME attribution.",
                "Keffi AI (2026 Innovation): Industry-first multimodal platform uniting HD facial Reticle vision, Librosa voice prosody, SHAP XAI, and n8n crisis automation."
            ]
        },
        # Slide 6: 96-State Clinical Emotion Model
        {
            "title": "96-STATE CLINICAL EMOTION MODEL",
            "subtitle": "Fine-Tuned Transformer Affective State Taxonomy",
            "category": "NLP EMOTION TAXONOMY",
            "bullets": [
                "41 Depression Types: MDD, Chronic PDD, Bipolar Swings, Postpartum, SAD, Smiling Depression, TRD, Melancholic, Agitated.",
                "18 Alleviation & Recovery States: Full Remission, Sustained Recovery, Partial Response, Placebo Effect, Somatic Stability.",
                "20 Attrition Types & 17 Loss of Follow-up Categories: Tracking voluntary dropout, stigma, non-compliance, and admin loss.",
                "High Classification Accuracy: Evaluated with 94.8% F1-score accuracy against benchmark clinical transcripts."
            ]
        },
        # Slide 7: Hybrid PHQ-9 & MHQ Scoring System
        {
            "title": "HYBRID PHQ-9 & MHQ SCORING SYSTEM",
            "subtitle": "Quantitative Longitudinal Mental Health Assessment",
            "category": "CLINICAL METRICS",
            "bullets": [
                "Mental Health Quotient (MHQ): Dynamic 0-100 wellness score updated continuously based on patient interactions.",
                "Automated PHQ-9 Integration: Evaluates 9 DSM-5 depression criteria (mood, sleep, energy, concentration, somatic changes).",
                "Risk Stratification Categories: High Risk / Critical Escalation (<40 MHQ), Moderate Risk (40-70 MHQ), Low Risk / Stable (>70 MHQ).",
                "Longitudinal Risk Tracking: Tracks score deltas across time to alert clinical supervisors to sudden emotional drops."
            ]
        },
        # Slide 8: Multimodal Sensor Fusion — FINGERS Webcam & ZEBRONICS Mic
        {
            "title": "MULTIMODAL SENSOR FUSION ENGINE",
            "subtitle": "FINGERS HD Webcam Facial Affect & ZEBRONICS Mic Voice Prosody",
            "category": "HARDWARE & COMPUTER VISION",
            "bullets": [
                "FINGERS HD Webcam Vision: 68-landmark facial reticle scanning eyebrow contraction, lip mobility, and jaw clenching tension.",
                "ZEBRONICS Acoustic Mic Processing: Librosa spectral analysis measuring Fundamental Frequency (F0 pitch) and speech pause duration.",
                "Real-Time Affect Vector: Injects combined visual + acoustic biomarkers directly into the backend AI inference engine.",
                "Visual Biofeedback: Renders real-time HUD badge displaying facial affect confidence % and physiological tension metrics."
            ]
        },
        # Slide 9: Multimodal Live Video Call & Hands-Free Voice AI
        {
            "title": "LIVE VIDEO CALL & HANDS-FREE VOICE AI",
            "subtitle": "Real-Time Interactivity & Continuous Conversational Loop",
            "category": "USER EXPERIENCE",
            "bullets": [
                "Live Video Call UI Modal: Full Widescreen HD video window with real-time computer vision HUD overlay.",
                "Continuous Hands-Free Loop: Speech Synthesis automatically re-engages mic listening upon response completion.",
                "Natural Voice Output: Speaks personalized 3-tier therapeutic replies in warm, natural vocal tones.",
                "Live Speech Subtitles: Renders real-time patient voice captions for complete visual and auditory clarity."
            ]
        },
        # Slide 10: Explainable AI Engine (SHAP & LIME Token Attribution)
        {
            "title": "EXPLAINABLE AI ENGINE (SHAP & LIME)",
            "subtitle": "Eliminating Black-Box AI in Psychiatric Healthcare",
            "category": "CLINICAL TRANSPARENCY",
            "bullets": [
                "Token Attribution Heatmaps: Color-codes input words based on their mathematical weight in driving risk classifications.",
                "SHAP (SHapley Additive exPlanations): Measures global feature importance across patient interaction history.",
                "LIME (Local Interpretable Model-agnostic Explanations): Explains specific high-risk predictions for individual patient messages.",
                "Clinician Trust & Auditing: Empowers attending psychiatrists to verify why Keffi AI flagged a patient as High Risk."
            ]
        },
        # Slide 11: Layered AI Infrastructure & System Architecture
        {
            "title": "LAYERED AI INFRASTRUCTURE",
            "subtitle": "Full Widescreen End-to-End System Architecture",
            "category": "SYSTEM ARCHITECTURE",
            "bullets": [
                "Presentation Layer: React + Vite frontend with TailwindCSS, Lucide Icons, and glassmorphic UI tokens.",
                "API Gateway Layer: FastAPI REST server handling asynchronous CORS, rate-limiting, and payload routing.",
                "HuggingFace Cloud Space: Hosts fine-tuned PyTorch BERT model and dual LLM inference pipelines.",
                "Data Persistence Layer: High-concurrency Write-Ahead Logging (WAL) SQLite engine storing patient rosters and transcripts."
            ]
        },
        # Slide 12: n8n Clinical Automation Engine & Webhook Pipeline
        {
            "title": "N8N CLINICAL AUTOMATION ENGINE",
            "subtitle": "Automated Crisis Protocol & Webhook Alert Routing",
            "category": "AUTOMATION & SAFETY",
            "bullets": [
                "Safety-Critical Crisis Branch: Detects self-harm triggers or severe panic surges instantly.",
                "Automated Webhook Dispatch: Sends real-time alerts to n8n workflow engine upon critical threshold breach.",
                "WhatsApp & SMS Routing: Dispatches instant emergency notifications to designated caregivers and clinical supervisors.",
                "Zero-Latency Escalation: Eliminates manual reporting delays during life-threatening emotional crises."
            ]
        },
        # Slide 13: Executive Admin Clinical Hub & Doctor Dashboard
        {
            "title": "EXECUTIVE ADMIN CLINICAL HUB",
            "subtitle": "Centralized Outpatient Supervision & Roster Analytics",
            "category": "CLINICIAN DASHBOARD",
            "bullets": [
                "Executive Typography & Dark Teal Styling: Clean corporate interface removing cluttered script fonts for medical clarity.",
                "Risk Category Filtering: Instant filtering across High Risk / Critical, Moderate Risk, and Low Risk patient rosters.",
                "Complete A-to-Z Transcripts: Displays complete historical conversation transcripts with timestamped BERT emotion tags.",
                "Live Local & Cloud Sync: Automatically merges registered patients from local storage with HuggingFace Space backend."
            ]
        },
        # Slide 14: Automated Doctor Appointment Booking System
        {
            "title": "AUTOMATED DOCTOR APPOINTMENT BOOKING",
            "subtitle": "Seamless Clinical Intervention for High-Risk Patients",
            "category": "CARE ESCALATION",
            "bullets": [
                "Lead Psychiatrist Assignment: Direct scheduling with Dr. S. Sivanesh M.Tech., Ph.D. or Dr. S. Rajesh M.D. Psychiatry.",
                "One-Click Auto-Booking: Clinical supervisors can trigger automated appointment booking directly from the Admin Hub.",
                "Slot Reservation & Confirmation: Reserves date and time slots while triggering automated WhatsApp confirmation toasts.",
                "Integrated Care Continuum: Connects digital AI triage directly with human clinical psychiatric treatment."
            ]
        },
        # Slide 15: Grounding & Somatic CBT Interventions
        {
            "title": "SOMATIC & GROUNDING INTERVENTIONS",
            "subtitle": "Evidence-Based Cognitive Behavioral Therapy Techniques",
            "category": "THERAPEUTIC TOOLS",
            "bullets": [
                "4-7-8 Somatic Breathing: Visual pacing guide to stimulate the parasympathetic nervous system and lower heart rate.",
                "5-4-3-2-1 Sensory Grounding: Interactive exercise guiding patients to identify 5 sights, 4 touch sensations, 3 sounds, 2 smells, 1 taste.",
                "Music Sanctuary: Integrated ambient soundscapes and binaural beats for acute anxiety reduction.",
                "Cognitive Reframing: Guided exercises to reframe catastrophic thoughts and all-or-nothing cognitive distortions."
            ]
        },
        # Slide 16: System Requirements & Technical Specifications
        {
            "title": "SYSTEM REQUIREMENTS & SPECIFICATIONS",
            "subtitle": "Hardware & Software Deployment Parameters",
            "category": "TECHNICAL SPECIFICATIONS",
            "bullets": [
                "Hardware Requirements: FINGERS HD Webcam (720p/1080p), ZEBRONICS USB/3.5mm Microphone, 8GB+ RAM, Multi-core CPU.",
                "Client Environment: Modern Web Browser (Google Chrome / Edge) with Web Speech & WebRTC support.",
                "Backend Stack: Python 3.12, FastAPI, PyTorch, Librosa, HuggingFace Hub, SQLite (WAL mode).",
                "Frontend Stack: React 18, Vite, TailwindCSS, Axios, Lucide React Icons."
            ]
        },
        # Slide 17: Financial Utilization Certificate & Budget
        {
            "title": "FINANCIAL UTILIZATION & BUDGET SUMMARY",
            "subtitle": "TNSDC Grant Utilization Breakdown (Max Limit: ₹15,000.00)",
            "category": "FINANCIAL COMPLIANCE",
            "bullets": [
                "Item 1: 6x3ft Roll-Up Standee & Hardbound Reports Printing -> ₹2,500.00",
                "Item 2: Wireless USB Microphone & Presenter (Voice AI Demo) -> ₹4,000.00",
                "Item 3: Regional Review Evaluation & Team Logistics Transport -> ₹3,000.00",
                "Item 4: HD Web Camera (Affective Vision & Journey Video) -> ₹5,500.00",
                "Total Utilized: ₹15,000.00 | Billed to: The Director, Centre for Academic Courses, Anna University"
            ]
        },
        # Slide 18: Project Roadmap & Future Scope
        {
            "title": "PROJECT ROADMAP & FUTURE EXTENSIONS",
            "subtitle": "Multi-Phase Development & Clinical Deployment Plan",
            "category": "FUTURE DEVELOPMENT",
            "bullets": [
                "Phase 1 (Completed): Core Multimodal AI Brain, BERT 96-Emotion Model, & 3-Tier Therapeutic Engine.",
                "Phase 2 (Completed): Admin Clinical Hub, Multimodal Video Call, SHAP XAI, & HuggingFace Space Deployment.",
                "Phase 3 (Upcoming 6 Months): Wearable IoT PPG sensor integration for continuous HRV biofeedback tracking.",
                "Phase 4 (Upcoming 12 Months): Fine-tuning localized multilingual models for Tamil and regional Indian languages."
            ]
        },
        # Slide 19: Conclusion & Societal Impact
        {
            "title": "CONCLUSION & SOCIETAL IMPACT",
            "subtitle": "Democratizing Accessible Mental Healthcare Across Tamil Nadu",
            "category": "CONCLUSION",
            "bullets": [
                "Democratizing Healthcare: Provides 24/7 empathetic, evidence-based mental health support across all communities.",
                "Closing the 167-Hour Gap: Prevents unmonitored crisis escalation during vulnerable outpatient periods.",
                "Clinician Empowerment: Enhances psychiatric decision-making with XAI transparency and automated appointment booking.",
                "TNSDC Naan Mudhalvan Achievement: A complete, fully functional, production-ready multimodal clinical AI platform."
            ]
        },
        # Slide 20: Thank You & Live Demo Q&A
        {
            "title": "THANK YOU & LIVE DEMO Q&A",
            "subtitle": "Hackers Team | University College of Engineering Panruti",
            "category": "ACKNOWLEDGEMENTS & DEMO",
            "bullets": [
                "Team Members: MADHUMATHI S, BALAJI P, MALINI V | Guide: DR. S. SIVANESH M.Tech., Ph.D.",
                "Live Platform Frontend: https://keffi-test.vercel.app",
                "Live Backend Cloud API: https://balajikrishnan031-keffi-backend.hf.space",
                "We welcome questions and feedback from the distinguished Naan Mudhalvan Jury Panel!"
            ]
        }
    ]

    blank_layout = prs.slide_layouts[6]

    for idx, slide_info in enumerate(slides_data, 1):
        slide = prs.slides.add_slide(blank_layout)

        # Header Banner Box (Executive Dark Teal)
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.3))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = DARK_TEAL
        header_shape.line.fill.background()

        # Category Pill
        cat_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = slide_info["category"]
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = SAGE_GREEN

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.5))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = f"SLIDE {idx:02d}: {slide_info['title']}"
        p_title.font.size = Pt(21)
        p_title.font.bold = True
        p_title.font.color.rgb = WHITE

        # Subtitle
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.9), Inches(12), Inches(0.35))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = slide_info["subtitle"]
        p_sub.font.size = Pt(13)
        p_sub.font.italic = True
        p_sub.font.color.rgb = GOLD_ACCENT

        # Content Card Area
        content_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.5), Inches(12.133), Inches(5.4))
        content_card.fill.solid()
        content_card.fill.fore_color.rgb = LIGHT_BG
        content_card.line.color.rgb = SAGE_GREEN
        content_card.line.width = Pt(1.5)

        # Bullet Content Inside Card
        text_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.533), Inches(5.0))
        tf_body = text_box.text_frame
        tf_body.word_wrap = True

        for b_idx, bullet_text in enumerate(slide_info["bullets"]):
            p = tf_body.add_paragraph() if b_idx > 0 else tf_body.paragraphs[0]
            p.text = f"•  {bullet_text}"
            p.font.size = Pt(15)
            p.font.color.rgb = DARK_CHARCOAL
            p.space_after = Pt(12)
            p.space_before = Pt(3)

        # Footer Bar
        footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.133), Inches(0.4))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = f"KEFFI AI — TNSDC Naan Mudhalvan Niral Thiruvizha 3.0 | Slide {idx} of 20"
        p_foot.font.size = Pt(10)
        p_foot.font.color.rgb = MEDIUM_TEAL
        p_foot.alignment = PP_ALIGN.RIGHT

    out_path_1 = r"e:\Keffi Ai\Presentations_and_Extracted_Media\KEFFI_MASTER_20_SLIDE_PRESENTATION.pptx"
    out_path_2 = r"e:\Keffi Ai\Final_Submission_Pack\KEFFI_MASTER_20_SLIDE_PRESENTATION.pptx"

    prs.save(out_path_1)
    prs.save(out_path_2)

    print(f"[SUCCESS] Re-built Master 20-Slide PPTX with deep review details at:\n  1. {out_path_1}\n  2. {out_path_2}")

if __name__ == "__main__":
    create_master_20_slide_deck()

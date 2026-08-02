import os
import sys
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import A3, landscape
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, HRFlowable
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY

def build_expanded_a3_poster():
    print("=== BUILDING EXPANDED A3 TIMES NEW ROMAN POSTER (LARGE FONTS & DETAILED TOPICS) ===")
    
    media_dir = r'e:\Keffi Ai\Presentations_and_Extracted_Media'
    pack_dir = r'e:\Keffi Ai\Final_Submission_Pack'

    target_pdf_main = os.path.join(media_dir, 'KEFFI_POSTER.pdf')
    target_pdf_final = os.path.join(media_dir, 'KEFFI_POSTER_FINAL.pdf')
    submission_pdf = os.path.join(pack_dir, '3_Project_Poster.pdf')

    target_pptx_main = os.path.join(media_dir, 'KEFFI_POSTER.pptx')
    submission_pptx = os.path.join(pack_dir, '3_Project_Poster.pptx')

    img_logo_left = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_1.png'
    img_logo_right = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_3.jpg'
    img_user_landing = r'e:\Keffi Ai\Documentation\new_user_landing_screenshot_v2.png'

    # -------------------------------------------------------------------------
    # PART 1: EXPANDED A3 LANDSCAPE PDF (1190.55 x 841.89 pt)
    # -------------------------------------------------------------------------
    doc = SimpleDocTemplate(
        target_pdf_final,
        pagesize=landscape(A3),
        leftMargin=0.30 * inch,
        rightMargin=0.30 * inch,
        topMargin=0.30 * inch,
        bottomMargin=0.30 * inch
    )

    styles = getSampleStyleSheet()

    PRIMARY_COLOR = colors.HexColor('#0F3838')   # Deep forest green
    SECONDARY_COLOR = colors.HexColor('#0D5050') # Heading color
    TEXT_DARK = colors.HexColor('#1E293B')        # Main body slate 800
    DIVIDER_COLOR = colors.HexColor('#0D7070')    # Section line divider

    style_title = ParagraphStyle(
        'TimesPosterTitle',
        fontName='Times-Bold',
        fontSize=24,
        leading=28,
        alignment=TA_CENTER,
        textColor=PRIMARY_COLOR
    )

    style_head = ParagraphStyle(
        'TimesSectionHead',
        fontName='Times-Bold',
        fontSize=15.0,
        leading=18.0,
        alignment=TA_LEFT,
        textColor=SECONDARY_COLOR,
        spaceAfter=3
    )

    style_body = ParagraphStyle(
        'TimesBody',
        fontName='Times-Roman',
        fontSize=10.5,
        leading=13.8,
        alignment=TA_JUSTIFY,
        textColor=TEXT_DARK,
        spaceAfter=6
    )

    style_bullet = ParagraphStyle(
        'TimesBullet',
        fontName='Times-Roman',
        fontSize=10.0,
        leading=13.2,
        alignment=TA_LEFT,
        textColor=TEXT_DARK,
        spaceAfter=4.0
    )

    story = []

    # TOP HEADER (CLEAN TEXT DIRECTLY ON WHITE BACKGROUND)
    header_html = """
    <font size="24" color="#0F3838"><b>KEFFI AI – A MENTAL HEALTH CHATBOT</b></font><br/>
    <font size="13" color="#0D5050"><b>UNIVERSITY COLLEGE OF ENGINEERING PANRUTI</b></font><br/>
    <font size="10.5" color="#334155"><i>(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering</i></font><br/>
    <font size="10" color="#1E293B"><b>TEAM NAME:</b> HACKERS TEAM &nbsp;&nbsp;|&nbsp;&nbsp; 
    <b>MEMBERS:</b> MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048)<br/>
    <b>GUIDE NAME:</b> DR. S. SIVANESH M.Tech., Ph.D. (Assistant Professor & Head of Department)</font>
    """
    p_header = Paragraph(header_html, style_title)

    img_l = Image(img_logo_left, width=85, height=85) if os.path.exists(img_logo_left) else Paragraph("", style_body)
    img_r = Image(img_logo_right, width=85, height=85) if os.path.exists(img_logo_right) else Paragraph("", style_body)

    header_table = Table([[img_l, p_header, img_r]], colWidths=[95, 955, 95])
    header_table.setStyle(TableStyle([
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 2),
        ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ('LEFTPADDING', (0,0), (-1,-1), 4),
        ('RIGHTPADDING', (0,0), (-1,-1), 4),
    ]))
    story.append(header_table)
    story.append(HRFlowable(width="100%", thickness=2.0, color=PRIMARY_COLOR, spaceBefore=4, spaceAfter=12))

    # Helper function for Times New Roman NO-BOX sections
    def make_times_section(heading_text, content_elements, width_pt):
        sec_story = []
        p_h = Paragraph(f"<b>{heading_text.upper()}</b>", style_head)
        sec_story.append(p_h)
        sec_story.append(HRFlowable(width="100%", thickness=1.5, color=DIVIDER_COLOR, spaceBefore=2, spaceAfter=6))
        for item in content_elements:
            sec_story.append(item)
        
        t = Table([[sec_story]], colWidths=[width_pt])
        t_style = [
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 2),
            ('RIGHTPADDING', (0,0), (-1,-1), 2),
            ('TOPPADDING', (0,0), (-1,-1), 2),
            ('BOTTOMPADDING', (0,0), (-1,-1), 6),
        ]
        t.setStyle(TableStyle(t_style))
        return t

    # COLUMN 1 (LEFT COLUMN - WIDTH: 375pt)
    c1_1 = Paragraph("Mental healthcare accessibility remains a critical global emergency due to exorbitant therapy costs, severe psychiatrist shortages, and persistent social stigma. Standard psychotherapy is restricted to 1 hour per week, leaving patients completely unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinically-grounded digital therapeutics platform engineered to bridge this gap, delivering 24/7 continuous affective monitoring, hands-free voice-to-voice therapeutics, and structured cognitive behavioral interventions.", style_body)
    sec_abstract = make_times_section("1. ABSTRACT", [c1_1], 375)

    c1_2 = [
        Paragraph("• <b>Fine-Tuned BERT 96-Emotion Engine:</b> Classifies complex, multi-layered user statements into 96 distinct clinical emotional categories with 94.2% top-3 accuracy.", style_bullet),
        Paragraph("• <b>3-Tier Clinical Response Framework:</b> Constructs therapeutic responses combining Rogerian empathetic validation, neurobiological psychoeducation, and CBT grounding skills.", style_bullet),
        Paragraph("• <b>Hands-Free Real-Time Voice AI:</b> Integrated Web Speech API audio processing allowing instant hands-free therapeutic conversation for patients in acute panic unable to type.", style_bullet),
        Paragraph("• <b>High-Concurrency WAL Database Engine:</b> Write-Ahead Logging database architecture eliminating concurrency lock crashes during peak user load.", style_bullet),
        Paragraph("• <b>Explainable AI (SHAP & LIME):</b> Visual feature attribution maps calculating exact token-level weights, empowering attending psychiatrists to audit AI decisions.", style_bullet),
        Paragraph("• <b>Admin Clinical Dashboard & Roster:</b> Real-time patient tracking displaying Days Inactive metrics and complete scrollable conversation transcripts.", style_bullet),
    ]
    sec_highlights = make_times_section("2. HIGHLIGHTS OF THE PROJECT", c1_2, 375)

    col1_sections = [sec_abstract, Spacer(1, 10), sec_highlights]

    # COLUMN 2 (CENTER COLUMN - WIDTH: 390pt)
    c2_1 = Paragraph("Keffi AI employs a multi-modal affective computing pipeline. User inputs (text utterances or real-time voice audio) pass through a dual-model processing cascade. The fine-tuned BERT transformer extracts high-dimensional semantic emotion vectors, while an integrated Librosa audio prosody analyzer extracts acoustic biomarkers including Fundamental Frequency (F0 pitch), RMS energy, and speech pause duration. Chronological session context is maintained by querying a WAL-enabled relational database and vector memory. If crisis signals or suicidal ideation are detected, automated n8n webhook workflows trigger immediate emergency contacts and clinical alerts.", style_body)
    sec_methodology = make_times_section("3. METHODOLOGY", [c2_1], 390)

    c2_2 = [
        Paragraph("<b>Step 1: Multimodal Data Ingestion:</b> Capture user text or real-time voice audio via Web Speech API endpoints.", style_bullet),
        Paragraph("<b>Step 2: BERT Emotion Classification:</b> Tokenize text input and classify across 96 fine-grained clinical emotional state vectors.", style_bullet),
        Paragraph("<b>Step 3: Isolated Context Retrieval:</b> Query WAL relational database and vector memory using patient ID to restore chat history.", style_bullet),
        Paragraph("<b>Step 4: 3-Tier Therapeutic Response Synthesis:</b> Synthesize Rogerian validation, neurobiological explanations (amygdala/cortisol surges), and CBT grounding exercises.", style_bullet),
        Paragraph("<b>Step 5: Acoustic Prosody & Risk Triage:</b> Compute pitch variance deltas; trigger automated n8n crisis alerts if suicidal ideation is flagged.", style_bullet),
        Paragraph("<b>Step 6: Clinician Dashboard Synchronization:</b> Update Mental Health Quotient (MHQ) telemetry and log complete transcripts in Admin Hub.", style_bullet),
    ]
    sec_procedure = make_times_section("4. STEP BY STEP PROCEDURE OR ALGORITHM", c2_2, 390)

    c2_3 = [
        Paragraph("• <b>GoEmotions Multi-Label Dataset:</b> 58,000 carefully curated Reddit comments annotated across 27 baseline emotion categories, extended to 96 clinical psychological states for transformer fine-tuning.", style_bullet),
        Paragraph("• <b>DAIC-WOZ Depression Audio Corpus:</b> Clinical interviews containing acoustic voice prosody benchmarks used to calibrate pitch (F0) variability and speech pause indicators.", style_bullet),
    ]
    sec_dataset = make_times_section("5. DATASET USED IF ANY", c2_3, 390)

    col2_sections = [sec_methodology, Spacer(1, 8), sec_procedure, Spacer(1, 8), sec_dataset]

    # COLUMN 3 (RIGHT COLUMN - WIDTH: 375pt)
    c3_1 = [
        Paragraph("<b>User Input:</b><br/>Spoken voice audio or written text statement:<br/><i>\"I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.\"</i>", style_bullet),
        Paragraph("<b>System Output:</b><br/>• <b>Tier 1 Validation:</b> <i>\"I hear how much pressure you're under. It's valid to feel overwhelmed.\"</i><br/>• <b>Tier 2 Psychoeducation:</b> <i>\"Your brain's amygdala is triggering a surge in cortisol and adrenaline.\"</i><br/>• <b>Tier 3 CBT Skill:</b> <i>\"Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.\"</i><br/>• <b>Voice Readout & Chips:</b> Spoken therapeutic readout and 3 interactive grounding chips.", style_bullet),
    ]
    sec_io = make_times_section("6. INPUT AND OUTPUT", c3_1, 375)

    c3_2 = Paragraph("Keffi AI successfully validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital mental health platform. By closing the 167-hour weekly care gap, Keffi AI empowers self-guided patient recovery while providing attending psychiatrists with real-time risk visibility.", style_body)
    sec_conclusion = make_times_section("7. CONCLUSION", [c3_2], 375)

    # USER'S NEW CRISP LIVE LANDING SCREENSHOT EMBEDDED HERE
    sec_img_element = Image(img_user_landing, width=365, height=195) if os.path.exists(img_user_landing) else Paragraph("", style_body)
    sec_system_img = make_times_section("DEVELOPED LIVE SYSTEM INTERFACE", [sec_img_element], 375)

    col3_sections = [sec_io, Spacer(1, 6), sec_conclusion, Spacer(1, 6), sec_system_img]

    main_grid = Table([[col1_sections, col2_sections, col3_sections]], colWidths=[375, 395, 375])
    main_grid.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 0),
        ('RIGHTPADDING', (0,0), (-1,-1), 0),
        ('TOPPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
    ]))

    story.append(main_grid)
    doc.build(story)
    print(f"[SUCCESS] Expanded A3 Times New Roman PDF Generated at: {target_pdf_final}")

    shutil.copyfile(target_pdf_final, target_pdf_main)
    shutil.copyfile(target_pdf_final, submission_pdf)

    # -------------------------------------------------------------------------
    # PART 2: EXPANDED A3 TIMES NEW ROMAN PPTX (KEFFI_POSTER.pptx)
    # -------------------------------------------------------------------------
    prs = Presentation()
    prs.slide_width = Inches(16.54)   # Exact A3 Landscape Slide Width
    prs.slide_height = Inches(11.69)  # Exact A3 Landscape Slide Height
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    COLOR_DARK_GREEN_PPT = RGBColor(15, 56, 56)      # #0F3838
    COLOR_TEAL_HEAD_PPT = RGBColor(13, 80, 80)       # #0D5050
    COLOR_SLATE_TEXT_PPT = RGBColor(30, 41, 59)      # #1E293B

    # Top Header Textbox
    tx_head = slide.shapes.add_textbox(Inches(1.5), Inches(0.2), Inches(13.54), Inches(1.5))
    tf_h = tx_head.text_frame
    tf_h.word_wrap = True

    p0 = tf_h.paragraphs[0]
    p0.text = "KEFFI AI – A MENTAL HEALTH CHATBOT"
    p0.font.name = "Times New Roman"
    p0.font.size = Pt(26)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_DARK_GREEN_PPT
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf_h.add_paragraph()
    p1.text = "UNIVERSITY COLLEGE OF ENGINEERING PANRUTI"
    p1.font.name = "Times New Roman"
    p1.font.size = Pt(14)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEAL_HEAD_PPT
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf_h.add_paragraph()
    p2.text = "(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering"
    p2.font.name = "Times New Roman"
    p2.font.size = Pt(11)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(51, 65, 85)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf_h.add_paragraph()
    p3.text = "TEAM NAME: HACKERS TEAM   |   MEMBERS: MADHUMATHI S, BALAJI P, MALINI V   |   GUIDE: DR. S. SIVANESH M.Tech., Ph.D."
    p3.font.name = "Times New Roman"
    p3.font.size = Pt(10.5)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_SLATE_TEXT_PPT
    p3.alignment = PP_ALIGN.CENTER

    if os.path.exists(img_logo_left):
        slide.shapes.add_picture(img_logo_left, Inches(0.4), Inches(0.2), width=Inches(1.3), height=Inches(1.3))
    if os.path.exists(img_logo_right):
        slide.shapes.add_picture(img_logo_right, Inches(14.8), Inches(0.2), width=Inches(1.3), height=Inches(1.3))

    # Divider line under header
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.75), Inches(15.74), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_DARK_GREEN_PPT
    line.line.color.rgb = COLOR_DARK_GREEN_PPT

    # Helper function for Times New Roman PPT sections
    def add_times_ppt_section(slide, heading_text, lines, left_in, top_in, width_in, height_in):
        tx_box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
        tf_box = tx_box.text_frame
        tf_box.word_wrap = True
        
        ph = tf_box.paragraphs[0]
        ph.text = heading_text.upper()
        ph.font.name = "Times New Roman"
        ph.font.size = Pt(15)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_TEAL_HEAD_PPT
        ph.alignment = PP_ALIGN.LEFT

        for line in lines:
            pl = tf_box.add_paragraph()
            pl.text = line
            pl.font.name = "Times New Roman"
            pl.font.size = Pt(10.5)
            pl.font.color.rgb = COLOR_SLATE_TEXT_PPT
            pl.space_after = Pt(4)
            pl.alignment = PP_ALIGN.LEFT

    # Column 1 (Left: 0.4 to 5.4 in)
    add_times_ppt_section(
        slide,
        "1. ABSTRACT",
        ["Mental healthcare accessibility remains a critical global emergency due to exorbitant therapy costs, severe psychiatrist shortages, and persistent social stigma. Standard psychotherapy is restricted to 1 hour per week, leaving patients completely unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinically-grounded digital therapeutics platform engineered to bridge this gap, delivering 24/7 continuous affective monitoring, hands-free voice-to-voice therapeutics, and structured cognitive behavioral interventions."],
        0.4, 1.9, 5.0, 2.5
    )

    add_times_ppt_section(
        slide,
        "2. HIGHLIGHTS OF THE PROJECT",
        [
            "• Fine-Tuned BERT 96-Emotion Engine: Classifies complex user statements into 96 distinct clinical emotional categories with 94.2% top-3 accuracy.",
            "• 3-Tier Clinical Response Framework: Constructs responses combining Rogerian validation, neurobiology, and CBT grounding skills.",
            "• Hands-Free Real-Time Voice AI: Integrated Web Speech API audio processing for patients in acute panic unable to type.",
            "• High-Concurrency WAL Database Engine: Write-Ahead Logging database architecture eliminating concurrency lock crashes.",
            "• Explainable AI (SHAP & LIME): Visual feature attribution maps calculating exact token-level weights for psychiatrist auditing.",
            "• Admin Clinical Dashboard & Roster: Displays Days Inactive metrics and complete scrollable conversation transcripts."
        ],
        0.4, 4.6, 5.0, 6.5
    )

    # Column 2 (Center: 5.7 to 10.9 in)
    add_times_ppt_section(
        slide,
        "3. METHODOLOGY",
        ["Keffi AI employs a multi-modal affective computing pipeline. User inputs pass through a dual-model cascade: fine-tuned BERT transformer classifies fine-grained emotion vectors, while an integrated Librosa audio prosody analyzer extracts acoustic biomarkers (F0 pitch, RMS energy, pause duration). Session context is maintained via a WAL-enabled relational DB and vector memory. If crisis signals occur, automated n8n workflows trigger emergency alerts immediately."],
        5.7, 1.9, 5.2, 2.5
    )

    add_times_ppt_section(
        slide,
        "4. STEP BY STEP PROCEDURE OR ALGORITHM",
        [
            "Step 1: Multimodal Data Ingestion: Capture user text utterance or real-time voice audio via Web Speech API.",
            "Step 2: BERT Emotion Classification: Tokenize text input and classify across 96 fine-grained clinical emotional state vectors.",
            "Step 3: Isolated Context Retrieval: Query WAL relational database and vector memory using patient ID.",
            "Step 4: 3-Tier Therapeutic Response Synthesis: Synthesize Rogerian validation, neurobiology, and CBT exercises.",
            "Step 5: Acoustic Prosody & Risk Triage: Compute pitch variance deltas; trigger automated n8n crisis alerts.",
            "Step 6: Clinician Dashboard Synchronization: Update Mental Health Quotient (MHQ) telemetry and log transcripts in Admin Hub."
        ],
        5.7, 4.6, 5.2, 4.1
    )

    add_times_ppt_section(
        slide,
        "5. DATASET USED IF ANY",
        [
            "• GoEmotions Multi-Label Dataset: 58,000 curated Reddit comments annotated across 27 baseline emotion categories, extended to 96 clinical psychological states for fine-tuning.",
            "• DAIC-WOZ Depression Audio Corpus: Clinical interviews containing acoustic voice prosody benchmarks used to calibrate pitch (F0) variability."
        ],
        5.7, 8.8, 5.2, 2.5
    )

    # Column 3 (Right: 11.2 to 16.1 in)
    add_times_ppt_section(
        slide,
        "6. INPUT AND OUTPUT",
        [
            "User Input: Spoken voice audio or written text statement ('I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.')",
            "System Output:",
            "• Tier 1 Validation: 'I hear how much pressure you're under. It's valid to feel overwhelmed.'",
            "• Tier 2 Psychoeducation: 'Your brain's amygdala is triggering a surge in cortisol.'",
            "• Tier 3 CBT Skill: 'Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.'",
            "• Voice Readout & Chips: Spoken therapeutic readout and interactive grounding chips."
        ],
        11.2, 1.9, 4.9, 3.5
    )

    add_times_ppt_section(
        slide,
        "7. CONCLUSION",
        ["Keffi AI successfully validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital mental health platform, closing the 167-hour weekly care gap."],
        11.2, 5.5, 4.9, 1.8
    )

    tx_arch = slide.shapes.add_textbox(Inches(11.2), Inches(7.4), Inches(4.9), Inches(0.4))
    tf_a = tx_arch.text_frame
    p_a = tf_a.paragraphs[0]
    p_a.text = "DEVELOPED LIVE SYSTEM INTERFACE"
    p_a.font.name = "Times New Roman"
    p_a.font.size = Pt(15)
    p_a.font.bold = True
    p_a.font.color.rgb = COLOR_TEAL_HEAD_PPT

    # USER'S NEW LANDING PAGE SCREENSHOT EMBEDDED ON PPT SLIDE
    if os.path.exists(img_user_landing):
        slide.shapes.add_picture(img_user_landing, Inches(11.2), Inches(7.9), width=Inches(4.9), height=Inches(3.4))

    prs.save(target_pptx_main)
    shutil.copyfile(target_pptx_main, submission_pptx)
    print(f"[SUCCESS] Expanded A3 Times New Roman PPTX Generated at: {target_pptx_main}")

if __name__ == "__main__":
    build_expanded_a3_poster()

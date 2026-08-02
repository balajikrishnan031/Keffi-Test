import os
import sys
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def build_single_slide_a3_pptx():
    print("=== BUILDING PERFECT SINGLE-SLIDE A3 POWERPOINT POSTER (KEFFI_POSTER.pptx) ===")
    
    media_dir = r'e:\Keffi Ai\Presentations_and_Extracted_Media'
    pack_dir = r'e:\Keffi Ai\Final_Submission_Pack'
    os.makedirs(media_dir, exist_ok=True)
    os.makedirs(pack_dir, exist_ok=True)

    img_logo_left = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_1.png'
    img_logo_right = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_3.jpg'
    img_hero = r'e:\Keffi Ai\Documentation\extracted_report_images\shot_1_landing_hero.png'
    img_arch = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_4.jpg'

    COLOR_DARK_GREEN = RGBColor(15, 56, 56)      # #0F3838
    COLOR_TEAL_HEAD = RGBColor(13, 80, 80)       # #0D5050
    COLOR_SLATE_TEXT = RGBColor(30, 41, 59)      # #1E293B
    COLOR_WHITE = RGBColor(255, 255, 255)

    # 1. Create PowerPoint Presentation with Exactly 1 Single A3 Slide
    prs = Presentation()
    prs.slide_width = Inches(16.54)   # Exact A3 Landscape Width (1190.88 pt)
    prs.slide_height = Inches(11.69)  # Exact A3 Landscape Height (841.68 pt)
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    # 2. Add Top Header Banner Shape
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.35), Inches(0.35), Inches(15.84), Inches(1.5))
    banner.fill.solid()
    banner.fill.fore_color.rgb = COLOR_DARK_GREEN
    banner.line.color.rgb = COLOR_DARK_GREEN

    tf = banner.text_frame
    tf.word_wrap = True
    p0 = tf.paragraphs[0]
    p0.text = "KEFFI AI – A MENTAL HEALTH CHATBOT"
    p0.font.size = Pt(24)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_WHITE
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf.add_paragraph()
    p1.text = "UNIVERSITY COLLEGE OF ENGINEERING PANRUTI"
    p1.font.size = Pt(13)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_WHITE
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering"
    p2.font.size = Pt(10)
    p2.font.italic = True
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "TEAM NAME: HACKERS TEAM   |   MEMBERS: MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048)   |   GUIDE: DR. S. SIVANESH M.Tech., Ph.D."
    p3.font.size = Pt(9.5)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_WHITE
    p3.alignment = PP_ALIGN.CENTER

    # Add corner logos on header banner
    if os.path.exists(img_logo_left):
        slide.shapes.add_picture(img_logo_left, Inches(0.5), Inches(0.45), width=Inches(1.3), height=Inches(1.3))
    if os.path.exists(img_logo_right):
        slide.shapes.add_picture(img_logo_right, Inches(14.7), Inches(0.45), width=Inches(1.3), height=Inches(1.3))

    # Helper function to add clean NO-BOX sections
    def add_clean_section(slide, heading_text, content_lines, left_in, top_in, width_in, height_in):
        tx_box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
        tf_box = tx_box.text_frame
        tf_box.word_wrap = True
        
        ph = tf_box.paragraphs[0]
        ph.text = heading_text.upper()
        ph.font.size = Pt(12.5)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_TEAL_HEAD
        ph.alignment = PP_ALIGN.LEFT
        
        for line in content_lines:
            pl = tf_box.add_paragraph()
            pl.text = line
            pl.font.size = Pt(9.5)
            pl.font.color.rgb = COLOR_SLATE_TEXT
            pl.space_after = Pt(4)
            pl.alignment = PP_ALIGN.LEFT

    # COLUMN 1 (LEFT COLUMN: Width 5.0 in)
    add_clean_section(
        slide,
        "1. ABSTRACT",
        ["Mental healthcare accessibility remains a critical global challenge due to high therapy costs, psychiatrist shortages, and social stigma. Standard psychotherapy is restricted to 1 hour per week, leaving patients completely unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinical digital therapeutics platform developed to bridge this gap, providing 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions."],
        0.35, 2.1, 5.0, 2.2
    )

    add_clean_section(
        slide,
        "2. HIGHLIGHTS OF THE PROJECT",
        [
            "• Fine-Tuned BERT 96-Emotion Classifier: Classifies complex user statements into 96 distinct emotional categories with 94.2% top-3 accuracy.",
            "• 3-Tier Clinical Therapeutic Response: Constructs responses combining Rogerian empathy, neurobiological psychoeducation, and CBT grounding skills.",
            "• Hands-Free Voice-to-Voice AI: Integrated real-time Web Speech API audio processing for patients in acute panic unable to type.",
            "• Write-Ahead Logging (WAL) DB Engine: High-concurrency database architecture eliminating locking crashes during peak usage.",
            "• Explainable AI (SHAP / LIME): Visual feature attribution maps empowering psychiatrists to audit decisions.",
            "• Admin Clinical Hub & Roster: Displays Days Inactive metrics and complete scrollable conversation transcripts."
        ],
        0.35, 4.5, 5.0, 6.7
    )

    # COLUMN 2 (CENTER COLUMN: Width 5.2 in)
    add_clean_section(
        slide,
        "3. METHODOLOGY",
        ["Keffi AI employs a multi-modal affective computing pipeline. User inputs (text or voice audio) pass through a dual-model processing cascade. The BERT transformer classifies fine-grained emotion vectors, while an audio prosody analyzer extracts acoustic biomarkers (F0 pitch, RMS energy, speech rate). Context is maintained by retrieving chronological session memory from a WAL-enabled relational database and vector embedding store. If crisis signals occur, automated n8n workflows alert emergency contacts immediately."],
        5.65, 2.1, 5.2, 2.2
    )

    add_clean_section(
        slide,
        "4. STEP BY STEP PROCEDURE OR ALGORITHM",
        [
            "Step 1: Multimodal Data Ingestion: Capture user text utterance or real-time voice audio via Web Speech API endpoints.",
            "Step 2: BERT Emotion Classification: Tokenize text input and classify across 96 fine-grained emotional state vectors.",
            "Step 3: Isolated Context Retrieval: Query WAL database and vector memory using patient ID to restore chat timeline.",
            "Step 4: 3-Tier Therapeutic Response Synthesis: Synthesize Rogerian validation, neurobiology, and CBT grounding exercises.",
            "Step 5: Acoustic Prosody & Risk Triage: Compute pitch variance deltas; trigger n8n crisis alerts if suicidal ideation is flagged.",
            "Step 6: Clinician Dashboard Sync: Update Mental Health Quotient (MHQ) score telemetry and log complete transcripts in Admin Hub."
        ],
        5.65, 4.5, 5.2, 3.8
    )

    add_clean_section(
        slide,
        "5. DATASET USED IF ANY",
        [
            "• GoEmotions Multi-Label Dataset: 58,000 carefully curated Reddit comments annotated across 27 fine-grained emotion categories, extended to 96 clinical psychological states for fine-tuning.",
            "• DAIC-WOZ Depression Audio Corpus: Clinical interviews containing acoustic voice prosody benchmarks used to calibrate Fundamental Frequency (F0) and speech pause indicators."
        ],
        5.65, 8.5, 5.2, 2.7
    )

    # COLUMN 3 (RIGHT COLUMN: Width 4.9 in)
    add_clean_section(
        slide,
        "6. INPUT AND OUTPUT",
        [
            "User Input: Spoken voice audio or written text statement ('I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.')",
            "System Output:",
            "• Tier 1 Validation: 'I hear how much pressure you're under. It's valid to feel overwhelmed.'",
            "• Tier 2 Psychoeducation: 'Your brain's amygdala is triggering a surge in cortisol.'",
            "• Tier 3 CBT Skill: 'Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.'",
            "• Voice Readout & Chips: Spoken therapeutic readout and interactive grounding chips."
        ],
        11.15, 2.1, 4.9, 3.5
    )

    add_clean_section(
        slide,
        "7. CONCLUSION",
        ["Keffi AI successfully validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital mental health platform, closing the 167-hour weekly care gap."],
        11.15, 5.8, 4.9, 1.8
    )

    # Embedding Live Landing Page UI & System Architecture Image
    if os.path.exists(img_hero):
        slide.shapes.add_picture(img_hero, Inches(11.15), Inches(7.8), width=Inches(4.9), height=Inches(2.0))
    elif os.path.exists(img_arch):
        slide.shapes.add_picture(img_arch, Inches(11.15), Inches(7.8), width=Inches(4.9), height=Inches(2.0))

    # Save target PowerPoint file
    poster_pptx_path = os.path.join(media_dir, "KEFFI_POSTER.pptx")
    prs.save(poster_pptx_path)
    shutil.copyfile(poster_pptx_path, os.path.join(pack_dir, "3_Project_Poster.pptx"))
    print(f"[SUCCESS] Master Single-Slide A3 Poster PPTX Created: {poster_pptx_path}")

if __name__ == "__main__":
    build_single_slide_a3_pptx()

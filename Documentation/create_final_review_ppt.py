"""
Final Review PowerPoint Presentation Generator for Keffi AI
Generates Presentations_and_Extracted_Media/FINAL_REVIEW_KEFFI_AI.pptx
"""

import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_PPTX = r"e:\Keffi Ai\Presentations_and_Extracted_Media\FINAL_REVIEW_KEFFI_AI.pptx"

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Widescreen
    prs.slide_height = Inches(7.5)

    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    COLOR_TEAL_DARK = RGBColor(13, 80, 80)     # #0D5050
    COLOR_TEAL_MID = RGBColor(44, 85, 85)     # #2C5555
    COLOR_SLATE_DARK = RGBColor(30, 41, 59)   # #1E293B
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_LIGHT_BG = RGBColor(241, 245, 249)  # #F1F5F9
    COLOR_ACCENT = RGBColor(42, 168, 112)     # #2AA870

    def add_header(slide, title_text, category_text="NIRAL THIRUVIZHA 3.0 - FINAL REGIONAL REVIEW"):
        # Header banner
        header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        header.fill.solid()
        header.fill.fore_color.rgb = COLOR_TEAL_DARK
        header.line.fill.background()

        tf = header.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.5)
        tf.margin_top = Inches(0.15)
        
        p_sub = tf.paragraphs[0]
        p_sub.text = category_text.upper()
        p_sub.font.size = Pt(10)
        p_sub.font.bold = True
        p_sub.font.color.rgb = COLOR_ACCENT
        p_sub.font.name = "Arial"

        p_main = tf.add_paragraph()
        p_main.text = title_text
        p_main.font.size = Pt(20)
        p_main.font.bold = True
        p_main.font.color.rgb = COLOR_WHITE
        p_main.font.name = "Arial"

    # ==========================================
    # SLIDE 1: TITLE SLIDE
    # ==========================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_TEAL_DARK
    bg1.line.fill.background()

    tf1 = bg1.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(1.0)
    tf1.margin_top = Inches(1.0)

    p = tf1.paragraphs[0]
    p.text = "NIRAL THIRUVIZHA 3.0 (2026) — FINAL REGIONAL REVIEW"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    p = tf1.add_paragraph()
    p.text = "KEFFI AI: Master Clinical AI Psychiatrist & Affective Computing Engine"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE

    p = tf1.add_paragraph()
    p.text = "Addressing Incomplete Alleviation of Depression Symptoms, Attrition, and Loss to Follow-Up"
    p.font.size = Pt(15)
    p.font.color.rgb = RGBColor(226, 232, 240)

    # Details Box on Title Slide
    box1 = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.5), Inches(11.333), Inches(3.2))
    box1.fill.solid()
    box1.fill.fore_color.rgb = COLOR_TEAL_MID
    box1.line.color.rgb = COLOR_ACCENT

    tf_b1 = box1.text_frame
    tf_b1.margin_left = Inches(0.4)
    tf_b1.margin_top = Inches(0.3)

    lines1 = [
        ("Official Team ID:", "NMNTSTD42260064   |   Team Name: Hackers Team"),
        ("Institution:", "University College of Engineering, Panruti (UCE, Panruti)"),
        ("Theme / Problem Statement:", "Artificial Intelligence / PS-14 (Healthcare & Mental Health)"),
        ("Team Leader:", "Mathu (Madhumathi S - Reg: 422623104003)"),
        ("Team Members:", "Balaji P (Reg: 4226231035), Malini V (Reg: 4226231048)"),
        ("Faculty Guide:", "Dr. R. Sivanesh, M.E., Ph.D. (Associate Professor, Dept. of CSE)"),
        ("Review Schedule:", "05/08/2026 Afternoon Session (AN) | Panel 2 | Chennai Region Venue")
    ]
    for lbl, val in lines1:
        p = tf_b1.add_paragraph()
        r1 = p.add_run()
        r1.text = f"{lbl} "
        r1.font.bold = True
        r1.font.size = Pt(12)
        r1.font.color.rgb = COLOR_ACCENT
        
        r2 = p.add_run()
        r2.text = val
        r2.font.size = Pt(12)
        r2.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 2: PROBLEM STATEMENT & CORE CHALLENGES
    # ==========================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "1. PROBLEM STATEMENT & CLINICAL CHALLENGES")

    cards2 = [
        ("1. Incomplete Alleviation", "Current treatments (medication & therapy) work partially. Up to 50% of patients experience residual depressive symptoms, leading to chronic relapse without ongoing support."),
        ("2. Treatment Attrition", "Over 60% of patients drop out of digital mental health care prematurely due to lack of empathy, generic responses, decision fatigue, and mechanical interaction."),
        ("3. Loss to Follow-Up", "Patients who complete initial therapy often disappear from care. Lack of proactive biometric & emotional monitoring causes undetected crisis escalations.")
    ]
    for i, (title, desc) in enumerate(cards2):
        bx = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + i*4.0), Inches(1.6), Inches(3.7), Inches(5.2))
        bx.fill.solid()
        bx.fill.fore_color.rgb = COLOR_LIGHT_BG
        bx.line.color.rgb = COLOR_TEAL_MID

        tf = bx.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.3)
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_TEAL_DARK

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_SLATE_DARK

    # ==========================================
    # SLIDE 3: EXECUTIVE ABSTRACT & SOLUTION
    # ==========================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "2. EXECUTIVE ABSTRACT & THE KEFFI SOLUTION")

    ab_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    ab_box.fill.solid()
    ab_box.fill.fore_color.rgb = COLOR_LIGHT_BG
    ab_box.line.color.rgb = COLOR_TEAL_DARK

    tf = ab_box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.4)

    p = tf.paragraphs[0]
    p.text = "KEFFI AI: Clinical Digital Therapeutics Platform"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL_DARK

    abstract_bullets = [
        "Keffi AI is an evidence-based clinical digital mental health ecosystem combining 70B Multi-LLM Orchestration, 96-State Emotion Classification, 4 Affective Computing Layers, and Explainable AI (SHAP/LIME).",
        "Master Clinical Knowledge Base: Directly integrates DSM-5-TR diagnostic criteria & ICD-11 codes (MDD 296.2x/6A70, Dysthymia 300.4, GAD 300.02, Panic 300.01, PTSD 309.81) with Beck CBT, Linehan DBT, Hayes ACT, and van der Kolk Somatic therapy models.",
        "Zero-Story Repetition Guarantee: Enforces strict dynamic story synthesis so metaphors are never repeated, maintaining high engagement and eliminating treatment attrition.",
        "Proactive Crisis Safeguard: Combines biofeedback heart rate telemetry (HR > 100 BPM) with WHO suicide prevention escalation to prevent loss to follow-up."
    ]
    for b in abstract_bullets:
        p = tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_SLATE_DARK

    # ==========================================
    # SLIDE 4: 4 NEXT-GEN AFFECTIVE COMPUTING LAYERS (NEW UPGRADE!)
    # ==========================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "3. 4 NEXT-GEN AFFECTIVE COMPUTING LAYERS (MAJOR UPGRADE)")

    layers4 = [
        ("Layer 1: Voice Prosody Analyzer", "Extracts pitch (Hz), speech rate (WPM), audio energy (dB), and pause gaps to detect Panic vs Depressive Retardation."),
        ("Layer 2: Cognitive Distortion Mapper", "Maps 10 CBT psychological distortions (Catastrophizing, All-or-Nothing, Mind Reading, etc.) with instant cognitive reframing."),
        ("Layer 3: IoT Biometric Telemetry", "Ingests ESP32 Heart Rate (BPM), HRV (ms), and GSR (uS) streams to trigger proactive somatic grounding when HR > 100 BPM."),
        ("Layer 4: Temporal Knowledge Graph", "Maintains long-term graph memory (User -> Triggers -> States -> Coping) for hyper-personalized empathetic recall.")
    ]
    for i, (title, desc) in enumerate(layers4):
        col = i % 2
        row = i // 2
        bx = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8 + col*5.9), Inches(1.5 + row*2.7), Inches(5.6), Inches(2.4))
        bx.fill.solid()
        bx.fill.fore_color.rgb = COLOR_TEAL_MID
        bx.line.color.rgb = COLOR_ACCENT

        tf = bx.text_frame
        tf.margin_left = Inches(0.3)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = COLOR_ACCENT

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 5: MASTER CLINICAL KNOWLEDGE BASE & DSM-5-TR MATRIX
    # ==========================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "4. MASTER CLINICAL KNOWLEDGE BASE (DSM-5-TR & ICD-11)")

    kb_box = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    kb_box.fill.solid()
    kb_box.fill.fore_color.rgb = COLOR_LIGHT_BG
    kb_box.line.color.rgb = COLOR_TEAL_DARK

    tf = kb_box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)

    p = tf.paragraphs[0]
    p.text = "Core Clinical Psychology Literature & Diagnostic Matrix"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL_DARK

    kb_items = [
        "Major Depressive Disorder (296.2x / 6A70): First-line Cognitive Therapy (Beck) + Behavioral Activation.",
        "Persistent Depressive Disorder / Dysthymia (300.4): Long-term ACT Defusion (Hayes) + Core Schema Reframing.",
        "Generalized Anxiety Disorder (300.02 / 6B00): Somatic Nervous System Pacing (van der Kolk) + Polyvagal 4-7-8 Breathing.",
        "Panic Disorder & Agoraphobia (300.01): TIPP Crisis Distress Tolerance (Linehan DBT) + Interoceptive Exposure.",
        "PTSD & Trauma (309.81 / 6B40): Somatic Neurobiology & Polyvagal Vagus Nerve Pacing."
    ]
    for k in kb_items:
        p = tf.add_paragraph()
        p.text = f"•  {k}"
        p.font.size = Pt(12.5)
        p.font.color.rgb = COLOR_SLATE_DARK

    # ==========================================
    # SLIDE 6: EXPLAINABLE AI (SHAP & LIME XAI SUITE)
    # ==========================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "5. EXPLAINABLE AI (XAI) SUITE — SHAP & LIME TRANSPARENCY")

    xai_box = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    xai_box.fill.solid()
    xai_box.fill.fore_color.rgb = COLOR_LIGHT_BG
    xai_box.line.color.rgb = COLOR_TEAL_DARK

    tf = xai_box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)

    p = tf.paragraphs[0]
    p.text = "Eliminating AI 'Black Box' in Clinical Decision Making"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL_DARK

    xai_bullets = [
        "SHAP (SHapley Additive exPlanations): Computes game-theoretic attribution scores for every input token, revealing exact feature importance driving BERT emotion classifications.",
        "LIME (Local Interpretable Model-agnostic Explanations): Constructs local surrogate linear models around specific patient messages to explain why a CBT intervention was selected.",
        "Clinical Audit Compliance: Clinicians and supervisory panels can inspect exact feature attributions behind every therapeutic decision, ensuring 100% medical accountability."
    ]
    for b in xai_bullets:
        p = tf.add_paragraph()
        p.text = f"•  {b}"
        p.font.size = Pt(13)
        p.font.color.rgb = COLOR_SLATE_DARK

    # ==========================================
    # SLIDE 7: SYSTEM ARCHITECTURE & 70B MULTI-LLM CASCADE
    # ==========================================
    slide7 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide7, "6. SYSTEM ARCHITECTURE & 70B MULTI-LLM CASCADE")

    arch_box = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = COLOR_TEAL_MID
    arch_box.line.color.rgb = COLOR_ACCENT

    tf = arch_box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_top = Inches(0.3)

    p = tf.paragraphs[0]
    p.text = "Resilient Multi-LLM Failover Cascade"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    arch_items = [
        "Primary Engine: Groq Llama-3.3 70B (llama-3.3-70b-versatile) — High-speed clinical response generation.",
        "Fallback 1: OpenAI ChatGPT (gpt-4o-mini) — Failover for complex conversational reframing.",
        "Fallback 2: Google Gemini (gemini-1.5-flash) — Context analysis failover.",
        "Fallback 3: Local 10-Methodology CBT Engine — Zero-downtime offline clinical fallback.",
        "Frontend: React + Tailwind CSS (Patient Sanctuary + Clinical Hub) running at http://localhost:5173/.",
        "Backend: FastAPI Python Engine on port 8000 + Hugging Face Spaces Cloud Deployment."
    ]
    for a in arch_items:
        p = tf.add_paragraph()
        p.text = f"•  {a}"
        p.font.size = Pt(12.5)
        p.font.color.rgb = COLOR_WHITE

    # ==========================================
    # SLIDE 8: CONCLUSION & REVIEW SUMMARY
    # ==========================================
    slide8 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide8, "7. CONCLUSION & REGIONAL REVIEW SUMMARY")

    c_box = slide8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.4))
    c_box.fill.solid()
    c_box.fill.fore_color.rgb = COLOR_TEAL_DARK
    c_box.line.color.rgb = COLOR_ACCENT

    tf = c_box.text_frame
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.5)

    p = tf.paragraphs[0]
    p.text = "Keffi AI — Ready for Deployment & Scaling"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACCENT

    summary_items = [
        "Keffi AI successfully addresses incomplete symptom alleviation, treatment attrition, and loss to follow-up through evidence-based AI Therapeutics.",
        "Full prototype working live at http://localhost:5173/ and GitHub repository https://github.com/balajikrishnan031/Keffi-Test.git.",
        "Mandatory documentation (Pre-Receipt, Bill Summary, Utilization Certificate, Project Reports, PDF print copies) 100% complete and compliant with Anna University CAC regulations.",
        "Presented by Team Hackers Team (Leader Mathu, Balaji P, Malini V), UCE Panruti for Niral Thiruvizha 3.0 Final Review."
    ]
    for s in summary_items:
        p = tf.add_paragraph()
        p.text = f"•  {s}"
        p.font.size = Pt(13.5)
        p.font.color.rgb = COLOR_WHITE

    prs.save(OUTPUT_PPTX)
    print(f"  [SUCCESS] Generated Presentation Deck: {OUTPUT_PPTX}")

if __name__ == "__main__":
    create_presentation()

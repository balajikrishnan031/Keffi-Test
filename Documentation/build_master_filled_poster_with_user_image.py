import os
import sys
import shutil
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from reportlab.lib.pagesizes import landscape
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image, HRFlowable
)
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY

def build_master_filled_poster():
    print("=== GENERATING MASTER FILLED POSTER WITH NEW USER LANDING SCREENSHOT (TIMES NEW ROMAN) ===")
    
    media_dir = r'e:\Keffi Ai\Presentations_and_Extracted_Media'
    pack_dir = r'e:\Keffi Ai\Final_Submission_Pack'

    target_pdf_main = os.path.join(media_dir, 'KEFFI_POSTER.pdf')
    target_pdf_final = os.path.join(media_dir, 'KEFFI_POSTER_FINAL.pdf')
    submission_pdf = os.path.join(pack_dir, '3_Project_Poster.pdf')

    target_pptx_main = os.path.join(media_dir, 'KEFFI_POSTER.pptx')
    submission_pptx = os.path.join(pack_dir, '3_Project_Poster.pptx')

    img_logo_left = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_1.png'
    img_logo_right = r'e:\Keffi Ai\Documentation\extracted_poster_images\poster_img_3.jpg'
    img_landing_user = r'e:\Keffi Ai\Documentation\new_user_landing_screenshot.png'  # User's new screenshot

    # -------------------------------------------------------------------------
    # PART 1: GENERATE FILLED NO-BOX PDF (960 x 540 pt)
    # -------------------------------------------------------------------------
    POSTER_WIDTH = 960
    POSTER_HEIGHT = 540

    doc = SimpleDocTemplate(
        target_pdf_final,
        pagesize=(POSTER_WIDTH, POSTER_HEIGHT),
        leftMargin=12,
        rightMargin=12,
        topMargin=8,
        bottomMargin=8
    )

    styles = getSampleStyleSheet()

    PRIMARY_COLOR = colors.HexColor('#0F3838')   # Deep forest green
    SECONDARY_COLOR = colors.HexColor('#0D5050') # Heading color
    TEXT_DARK = colors.HexColor('#1E293B')        # Main body slate 800
    DIVIDER_COLOR = colors.HexColor('#0D7070')    # Section line divider

    style_title = ParagraphStyle(
        'TimesPosterTitle',
        fontName='Times-Bold',
        fontSize=19,
        leading=22,
        alignment=TA_CENTER,
        textColor=PRIMARY_COLOR
    )

    style_head = ParagraphStyle(
        'TimesSectionHead',
        fontName='Times-Bold',
        fontSize=11.0,
        leading=13.5,
        alignment=TA_LEFT,
        textColor=SECONDARY_COLOR,
        spaceAfter=2
    )

    style_body = ParagraphStyle(
        'TimesBody',
        fontName='Times-Roman',
        fontSize=8.1,
        leading=10.6,
        alignment=TA_JUSTIFY,
        textColor=TEXT_DARK,
        spaceAfter=5
    )

    style_bullet = ParagraphStyle(
        'TimesBullet',
        fontName='Times-Roman',
        fontSize=7.9,
        leading=10.4,
        alignment=TA_LEFT,
        textColor=TEXT_DARK,
        spaceAfter=3.0
    )

    story = []

    # TOP HEADER (NO GREEN BOX BACKGROUND)
    header_html = """
    <font size="19" color="#0F3838"><b>KEFFI AI – A MENTAL HEALTH CHATBOT</b></font><br/>
    <font size="10.5" color="#0D5050"><b>UNIVERSITY COLLEGE OF ENGINEERING PANRUTI</b></font><br/>
    <font size="8.5" color="#334155"><i>(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering</i></font><br/>
    <font size="8" color="#1E293B"><b>TEAM:</b> HACKERS TEAM &nbsp;|&nbsp; 
    <b>MEMBERS:</b> MADHUMATHI S (422623104003), BALAJI P (422623104035), MALINI V (422623104048) &nbsp;|&nbsp; 
    <b>GUIDE:</b> DR. S. SIVANESH M.Tech., Ph.D.</font>
    """
    p_header = Paragraph(header_html, style_title)

    img_l = Image(img_logo_left, width=70, height=70) if os.path.exists(img_logo_left) else Paragraph("", style_body)
    img_r = Image(img_logo_right, width=70, height=70) if os.path.exists(img_logo_right) else Paragraph("", style_body)

    header_table = Table([[img_l, p_header, img_r]], colWidths=[75, 786, 75])
    header_table.setStyle(TableStyle([
        ('ALIGN', (0,0), (-1,-1), 'CENTER'),
        ('VALIGN', (0,0), (-1,-1), 'MIDDLE'),
        ('TOPPADDING', (0,0), (-1,-1), 2),
        ('BOTTOMPADDING', (0,0), (-1,-1), 3),
        ('LEFTPADDING', (0,0), (-1,-1), 4),
        ('RIGHTPADDING', (0,0), (-1,-1), 4),
    ]))
    story.append(header_table)
    story.append(HRFlowable(width="100%", thickness=1.5, color=PRIMARY_COLOR, spaceBefore=2, spaceAfter=8))

    # Helper function for Times New Roman NO-BOX sections
    def make_times_section(heading_text, content_elements, width_pt):
        sec_story = []
        p_h = Paragraph(f"<b>{heading_text.upper()}</b>", style_head)
        sec_story.append(p_h)
        sec_story.append(HRFlowable(width="100%", thickness=1.2, color=DIVIDER_COLOR, spaceBefore=1, spaceAfter=4))
        for item in content_elements:
            sec_story.append(item)
        
        t = Table([[sec_story]], colWidths=[width_pt])
        t_style = [
            ('VALIGN', (0,0), (-1,-1), 'TOP'),
            ('LEFTPADDING', (0,0), (-1,-1), 2),
            ('RIGHTPADDING', (0,0), (-1,-1), 2),
            ('TOPPADDING', (0,0), (-1,-1), 2),
            ('BOTTOMPADDING', (0,0), (-1,-1), 4),
        ]
        t.setStyle(TableStyle(t_style))
        return t

    # COLUMN 1 (LEFT: 304pt)
    c1_1 = Paragraph("Mental healthcare accessibility remains a major global challenge due to high therapy costs, psychiatrist shortages, and social stigma. Psychotherapy is restricted to 1 hour per week, leaving patients unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinical digital therapeutics platform developed to bridge this gap, providing 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions.", style_body)
    sec_abstract = make_times_section("1. ABSTRACT", [c1_1], 304)

    c1_2 = [
        Paragraph("• <b>BERT 96-Emotion Classifier:</b> Classifies complex user statements into 96 distinct emotional categories with 94.2% top-3 accuracy.", style_bullet),
        Paragraph("• <b>3-Tier Clinical Response:</b> Combines Rogerian empathy, neurobiological psychoeducation, and CBT grounding skills.", style_bullet),
        Paragraph("• <b>Hands-Free Voice AI:</b> Real-time Web Speech API audio processing for patients in acute panic unable to type.", style_bullet),
        Paragraph("• <b>Write-Ahead Logging (WAL) DB:</b> High-concurrency database eliminating locking crashes during peak usage.", style_bullet),
        Paragraph("• <b>SHAP / LIME Explainable AI:</b> Visual feature attribution maps empowering psychiatrists to audit decisions.", style_bullet),
        Paragraph("• <b>Admin Clinical Hub:</b> Displays Days Inactive metrics and complete scrollable conversation transcripts.", style_bullet),
    ]
    sec_highlights = make_times_section("2. HIGHLIGHTS OF THE PROJECT", c1_2, 304)

    c1_3 = Paragraph("Keffi AI employs a multi-modal affective computing pipeline. User inputs (text or voice) pass through a dual-model cascade: BERT classifies fine-grained emotion vectors while a Librosa audio prosody analyzer extracts acoustic biomarkers (F0 pitch, RMS energy). Context is maintained via WAL relational DB and vector memory. If crisis signals occur, automated n8n workflows alert emergency contacts immediately.", style_body)
    sec_methodology = make_times_section("3. METHODOLOGY", [c1_3], 304)

    col1_sections = [sec_abstract, Spacer(1, 4), sec_highlights, Spacer(1, 4), sec_methodology]

    # COLUMN 2 (CENTER: 318pt)
    c2_1 = [
        Paragraph("<b>Step 1: Multimodal Ingestion:</b> Capture user text or voice audio via Web Speech API endpoints.", style_bullet),
        Paragraph("<b>Step 2: BERT Emotion Classification:</b> Classify input across 96 fine-grained emotional state vectors.", style_bullet),
        Paragraph("<b>Step 3: Isolated Context Retrieval:</b> Query WAL database and vector memory using patient ID.", style_bullet),
        Paragraph("<b>Step 4: 3-Tier Therapeutic Generation:</b> Synthesize Rogerian care, neurobiology, and CBT skills.", style_bullet),
        Paragraph("<b>Step 5: Acoustic Prosody & Risk Triage:</b> Compute pitch deltas; trigger n8n crisis alerts if suicidal ideation is flagged.", style_bullet),
        Paragraph("<b>Step 6: Clinician Hub Sync:</b> Update Mental Health Quotient (MHQ) telemetry and log transcripts in Admin Hub.", style_bullet),
    ]
    sec_procedure = make_times_section("4. STEP BY STEP PROCEDURE OR ALGORITHM", c2_1, 318)

    c2_2 = [
        Paragraph("• <b>GoEmotions Multi-Label Dataset:</b> 58,000 curated Reddit comments annotated across 27 emotion categories, extended to 96 clinical psychological states for fine-tuning.", style_bullet),
        Paragraph("• <b>DAIC-WOZ Depression Audio Corpus:</b> Clinical interviews containing acoustic voice prosody benchmarks used to calibrate pitch (F0) and speech pause indicators.", style_bullet),
    ]
    sec_dataset = make_times_section("5. DATASET USED IF ANY", c2_2, 318)

    c2_3 = Paragraph("Keffi AI validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital platform, closing the 167-hour weekly care gap.", style_body)
    sec_conclusion = make_times_section("7. CONCLUSION", [c2_3], 318)

    col2_sections = [sec_procedure, Spacer(1, 4), sec_dataset, Spacer(1, 4), sec_conclusion]

    # COLUMN 3 (RIGHT: 304pt)
    c3_1 = [
        Paragraph("<b>User Input:</b><br/>Spoken audio or text statement:<br/><i>\"I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.\"</i>", style_bullet),
        Paragraph("<b>System Output:</b><br/>• <b>Tier 1 Validation:</b> <i>\"I hear how much pressure you're under. It's valid to feel overwhelmed.\"</i><br/>• <b>Tier 2 Psychoeducation:</b> <i>\"Your brain's amygdala is triggering a surge in cortisol.\"</i><br/>• <b>Tier 3 CBT Skill:</b> <i>\"Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.\"</i><br/>• <b>Voice Readout & Chips:</b> Spoken therapeutic readout and interactive grounding chips.", style_bullet),
    ]
    sec_io = make_times_section("6. INPUT AND OUTPUT", c3_1, 304)

    # Embedding the USER'S NEW LIVE LANDING SCREENSHOT
    sec_img_element = Image(img_landing_user, width=294, height=160) if os.path.exists(img_landing_user) else Paragraph("", style_body)
    sec_system_img = make_times_section("DEVELOPED LIVE SYSTEM INTERFACE", [sec_img_element], 304)

    col3_sections = [sec_io, Spacer(1, 4), sec_system_img]

    main_grid = Table([[col1_sections, col2_sections, col3_sections]], colWidths=[310, 324, 310])
    main_grid.setStyle(TableStyle([
        ('VALIGN', (0,0), (-1,-1), 'TOP'),
        ('LEFTPADDING', (0,0), (-1,-1), 2),
        ('RIGHTPADDING', (0,0), (-1,-1), 2),
        ('TOPPADDING', (0,0), (-1,-1), 0),
        ('BOTTOMPADDING', (0,0), (-1,-1), 0),
    ]))

    story.append(main_grid)
    doc.build(story)
    print(f"[SUCCESS] Master Filled Times New Roman PDF Generated at: {target_pdf_final}")

    shutil.copyfile(target_pdf_final, target_pdf_main)
    shutil.copyfile(target_pdf_final, submission_pdf)

    # -------------------------------------------------------------------------
    # PART 2: TIMES NEW ROMAN PPTX (KEFFI_POSTER.pptx) WITH NEW USER IMAGE
    # -------------------------------------------------------------------------
    prs = Presentation()
    prs.slide_width = Inches(13.33)  # 16:9 Aspect Ratio
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    COLOR_DARK_GREEN_PPT = RGBColor(15, 56, 56)      # #0F3838
    COLOR_TEAL_HEAD_PPT = RGBColor(13, 80, 80)       # #0D5050
    COLOR_SLATE_TEXT_PPT = RGBColor(30, 41, 59)      # #1E293B

    # Header Textbox
    tx_head = slide.shapes.add_textbox(Inches(1.4), Inches(0.12), Inches(10.53), Inches(1.2))
    tf_h = tx_head.text_frame
    tf_h.word_wrap = True

    p0 = tf_h.paragraphs[0]
    p0.text = "KEFFI AI – A MENTAL HEALTH CHATBOT"
    p0.font.name = "Times New Roman"
    p0.font.size = Pt(21)
    p0.font.bold = True
    p0.font.color.rgb = COLOR_DARK_GREEN_PPT
    p0.alignment = PP_ALIGN.CENTER

    p1 = tf_h.add_paragraph()
    p1.text = "UNIVERSITY COLLEGE OF ENGINEERING PANRUTI"
    p1.font.name = "Times New Roman"
    p1.font.size = Pt(11.5)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEAL_HEAD_PPT
    p1.alignment = PP_ALIGN.CENTER

    p2 = tf_h.add_paragraph()
    p2.text = "(A Constituent College of Anna University, Chennai) • Department of Computer Science and Engineering"
    p2.font.name = "Times New Roman"
    p2.font.size = Pt(9.2)
    p2.font.italic = True
    p2.font.color.rgb = RGBColor(51, 65, 85)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf_h.add_paragraph()
    p3.text = "TEAM: HACKERS TEAM   |   MEMBERS: MADHUMATHI S, BALAJI P, MALINI V   |   GUIDE: DR. S. SIVANESH M.Tech., Ph.D."
    p3.font.name = "Times New Roman"
    p3.font.size = Pt(8.8)
    p3.font.bold = True
    p3.font.color.rgb = COLOR_SLATE_TEXT_PPT
    p3.alignment = PP_ALIGN.CENTER

    if os.path.exists(img_logo_left):
        slide.shapes.add_picture(img_logo_left, Inches(0.3), Inches(0.18), width=Inches(1.05), height=Inches(1.05))
    if os.path.exists(img_logo_right):
        slide.shapes.add_picture(img_logo_right, Inches(12.0), Inches(0.18), width=Inches(1.05), height=Inches(1.05))

    # Divider line under header
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.36), Inches(12.73), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_DARK_GREEN_PPT
    line.line.color.rgb = COLOR_DARK_GREEN_PPT

    # Helper function for Times New Roman PPT sections
    def add_times_ppt_section(slide, heading_text, lines, left_in, top_in, width_in, height_in):
        tx_box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
        tf_box = tx_box.text_frame
        tf_box.word_wrap = True
        
        ph = tf_box.paragraphs[0]
        ph.text = heading_text.upper()
        ph.font.name = "Times New Roman"
        ph.font.size = Pt(11.5)
        ph.font.bold = True
        ph.font.color.rgb = COLOR_TEAL_HEAD_PPT
        ph.alignment = PP_ALIGN.LEFT

        for line in lines:
            pl = tf_box.add_paragraph()
            pl.text = line
            pl.font.name = "Times New Roman"
            pl.font.size = Pt(8.2)
            pl.font.color.rgb = COLOR_SLATE_TEXT_PPT
            pl.space_after = Pt(2.5)
            pl.alignment = PP_ALIGN.LEFT

    # Column 1 (Left)
    add_times_ppt_section(
        slide,
        "1. ABSTRACT",
        ["Mental healthcare accessibility remains a major global challenge due to high therapy costs, psychiatrist shortages, and social stigma. Psychotherapy is restricted to 1 hour per week, leaving patients unmonitored during the remaining 167 hours of weekly vulnerability. Keffi AI is a clinical digital therapeutics platform developed to bridge this gap, providing 24/7 continuous affective monitoring, hands-free voice interaction, and structured psychological interventions."],
        0.2, 1.44, 4.1, 1.5
    )

    add_times_ppt_section(
        slide,
        "2. HIGHLIGHTS OF THE PROJECT",
        [
            "• BERT 96-Emotion Classifier: Classifies complex user statements into 96 distinct emotional categories with 94.2% accuracy.",
            "• 3-Tier Clinical Response: Combines Rogerian empathy, neurobiological psychoeducation, and CBT grounding skills.",
            "• Hands-Free Voice AI: Real-time Web Speech API audio processing for patients in acute panic unable to type.",
            "• Write-Ahead Logging (WAL) DB: High-concurrency database eliminating locking crashes during peak usage.",
            "• SHAP / LIME Explainable AI: Visual feature attribution maps empowering psychiatrists to audit decisions.",
            "• Admin Clinical Hub: Displays Days Inactive metrics and complete scrollable conversation transcripts."
        ],
        0.2, 3.02, 4.1, 2.6
    )

    add_times_ppt_section(
        slide,
        "3. METHODOLOGY",
        ["Keffi AI employs a multi-modal affective computing pipeline. User inputs (text or voice) pass through a dual-model cascade: BERT classifies fine-grained emotion vectors while a Librosa audio prosody analyzer extracts acoustic biomarkers (F0 pitch, RMS energy). Context is maintained via WAL relational DB and vector memory. If crisis signals occur, automated n8n workflows alert emergency contacts immediately."],
        0.2, 5.68, 4.1, 1.6
    )

    # Column 2 (Center)
    add_times_ppt_section(
        slide,
        "4. STEP BY STEP PROCEDURE OR ALGORITHM",
        [
            "Step 1: Multimodal Ingestion: Capture user text or voice audio via Web Speech API endpoints.",
            "Step 2: BERT Emotion Classification: Classify input across 96 fine-grained emotional state vectors.",
            "Step 3: Isolated Context Retrieval: Query WAL database and vector memory using patient ID.",
            "Step 4: 3-Tier Therapeutic Generation: Synthesize Rogerian care, neurobiology, and CBT skills.",
            "Step 5: Acoustic Prosody & Risk Triage: Compute pitch deltas; trigger n8n crisis alerts if suicidal ideation is flagged.",
            "Step 6: Clinician Hub Sync: Update Mental Health Quotient (MHQ) telemetry and log transcripts in Admin Hub."
        ],
        4.5, 1.44, 4.3, 2.7
    )

    add_times_ppt_section(
        slide,
        "5. DATASET USED IF ANY",
        [
            "• GoEmotions Multi-Label Dataset: 58,000 curated Reddit comments annotated across 27 emotion categories, extended to 96 clinical psychological states for fine-tuning.",
            "• DAIC-WOZ Depression Audio Corpus: Clinical interviews containing acoustic voice prosody benchmarks used to calibrate pitch (F0) and speech pause indicators."
        ],
        4.5, 4.22, 4.3, 1.7
    )

    add_times_ppt_section(
        slide,
        "7. CONCLUSION",
        ["Keffi AI validates the integration of fine-grained emotion classification, high-concurrency database storage, hands-free voice therapeutics, and clinician tracking into a unified digital platform, closing the 167-hour weekly care gap."],
        4.5, 5.98, 4.3, 1.3
    )

    # Column 3 (Right)
    add_times_ppt_section(
        slide,
        "6. INPUT AND OUTPUT",
        [
            "User Input: Spoken audio or text statement ('I feel completely overwhelmed by exam deadlines, my heart is racing, and I can't sleep.')",
            "System Output:",
            "• Tier 1 Validation: 'I hear how much pressure you're under. It's valid to feel overwhelmed.'",
            "• Tier 2 Psychoeducation: 'Your brain's amygdala is triggering a surge in cortisol.'",
            "• Tier 3 CBT Skill: 'Try 4-7-8 breathing: Breathe in for 4s, hold for 7s, exhale for 8s.'",
            "• Voice Readout & Chips: Spoken therapeutic readout and interactive grounding chips."
        ],
        9.0, 1.44, 4.1, 2.8
    )

    tx_arch = slide.shapes.add_textbox(Inches(9.0), Inches(4.35), Inches(4.1), Inches(0.35))
    tf_a = tx_arch.text_frame
    p_a = tf_a.paragraphs[0]
    p_a.text = "DEVELOPED LIVE SYSTEM INTERFACE"
    p_a.font.name = "Times New Roman"
    p_a.font.size = Pt(11.5)
    p_a.font.bold = True
    p_a.font.color.rgb = COLOR_TEAL_HEAD_PPT

    # USER'S NEW LANDING PAGE SCREENSHOT EMBEDDED HERE
    if os.path.exists(img_landing_user):
        slide.shapes.add_picture(img_landing_user, Inches(9.0), Inches(4.75), width=Inches(4.1), height=Inches(2.55))

    prs.save(target_pptx_main)
    shutil.copyfile(target_pptx_main, submission_pptx)
    print(f"[SUCCESS] Master Filled Times New Roman PPTX Generated at: {target_pptx_main}")

if __name__ == "__main__":
    build_master_filled_poster()
